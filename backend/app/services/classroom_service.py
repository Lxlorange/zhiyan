from __future__ import annotations

from datetime import datetime
from pathlib import Path
import json
import re
import socket
import urllib.error
import urllib.request
from typing import Any, Optional

from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy import Select, select
from sqlalchemy.orm import Session, object_session, selectinload

from app.models.learning import (
    AgentTaskRecord,
    ClassroomResource,
    ClassroomSession,
    ClassroomSubmission,
    LearningProject,
    LearningProjectEvent,
    LearningSyllabusVersion,
    LearningSyllabusItem,
)
from app.models.user import User
from app.schemas import (
    ClassroomDialogueRequest,
    ClassroomNoteSaveRequest,
    ClassroomQuizSubmitRequest,
    ClassroomReflectionSubmitRequest,
    ClassroomSlidesCompleteRequest,
    ClassroomVisualizationGenerateRequest,
    ClassroomVoiceGenerateRequest,
    ProfileEntryUpdateRequest,
    SyllabusItemStatusRequest,
)
from app.core.config import get_settings
from app.services.knowledge_ingestion_service import build_rag_context
from app.services.json_repair_service import LLMJsonParseError, parse_llm_json
from app.services.llm_client import LLMConfigurationError, LLMResponseError, qwen_chat_json, resolve_chat_config, validate_qwen_config
from app.services.syllabus_service import update_syllabus_item_status
from app.services.adaptive_visualization_renderer import render_adaptive_visualization_html
from app.services.visualization_3d_renderer import render_three_physics_html
from app.services.formula_guidance import FORMULA_OUTPUT_INSTRUCTIONS


PREGENERATE_CLASSROOM_LIMIT = 2
CLASSROOM_PACKAGE_MAX_TOKENS = 5200
CLASSROOM_REPAIR_MAX_TOKENS = 4200


class _SlideSpec(BaseModel):
    title: str
    layout: str
    bullets: list[str] = Field(min_length=2, max_length=6)
    speaker_notes: str
    visual_hint: str = ""
    visual_blocks: list[dict[str, Any]] = Field(min_length=2, max_length=6)
    side_panel: dict[str, Any]
    takeaways: list[str] = Field(min_length=2, max_length=5)
    source_refs: list[dict[str, str]] = Field(min_length=1, max_length=3)
    example: str = ""
    misconception: str = ""
    interaction_prompt: str = ""


class _QuizSpec(BaseModel):
    id: str
    prompt: str
    question_type: str = "single"
    options: list[dict[str, str]] = Field(min_length=3, max_length=6)
    answer: str
    explanation: str
    hint: str = ""
    difficulty: str = "medium"
    knowledge_point: str = ""


class _ConceptCardSpec(BaseModel):
    name: str
    explanation: str
    scenario: str
    misconception: str
    relation_to_project: str


class _DiagramSpec(BaseModel):
    title: str
    diagram_type: str
    mermaid: str
    explanation: str


class _VoiceScriptSpec(BaseModel):
    one_minute: str
    five_minutes: str
    segments: list[str] = Field(min_length=1, max_length=8)


class _ReproductionDemoSpec(BaseModel):
    title: str
    task: str
    input_format: str
    code_skeleton: str = ""
    steps: list[str] = Field(min_length=1, max_length=8)
    expected_output: str
    parameters: list[str] = Field(min_length=1, max_length=8)
    common_errors: list[str] = Field(min_length=1, max_length=8)
    report_suggestions: list[str] = Field(min_length=1, max_length=8)


class _GuidingQuestionSpec(BaseModel):
    prompt: str
    intent: str
    hint: str


class _ReadingSpec(BaseModel):
    title: str
    why: str
    source: str
    keywords: list[str] = Field(min_length=1, max_length=8)


class _ClassroomPackage(BaseModel):
    title: str
    learning_summary: str
    slides: list[_SlideSpec] = Field(min_length=4, max_length=6)
    concept_cards: list[_ConceptCardSpec] = Field(min_length=3, max_length=6)
    diagram: _DiagramSpec
    guiding_questions: list[_GuidingQuestionSpec] = Field(min_length=3, max_length=6)
    voice_script: _VoiceScriptSpec
    reproduction_demo: _ReproductionDemoSpec
    readings: list[_ReadingSpec] = Field(min_length=2, max_length=5)
    quiz: list[_QuizSpec] = Field(min_length=2, max_length=5)
    reflection_prompts: list[str] = Field(min_length=3, max_length=6)
    safety_notes: list[str]


class _VisualizationFrame(BaseModel):
    label: str
    metrics: dict[str, float]
    narrative: str
    active_nodes: list[str] = Field(default_factory=list)
    active_edges: list[str] = Field(default_factory=list)
    annotations: list[str] = Field(default_factory=list)
    camera_position: Optional[list[float]] = None
    camera_target: Optional[list[float]] = None


class _VisualizationControl(BaseModel):
    name: str
    label: str
    min_value: float
    max_value: float
    default_value: float
    description: str


class _PhysicsObjectSpec(BaseModel):
    id: str
    label: str
    role: str
    shape: str
    size: list[float] = Field(min_length=3, max_length=3)
    position: list[float] = Field(min_length=3, max_length=3)
    velocity: list[float] = Field(min_length=3, max_length=3)
    mass: float
    color: str
    particle_emitter: bool = False


class _PhysicsSceneSpec(BaseModel):
    scene_kind: str
    gravity: list[float] = Field(min_length=3, max_length=3)
    camera: dict[str, list[float]]
    objects: list[_PhysicsObjectSpec] = Field(min_length=4, max_length=9)
    annotations: list[dict[str, Any]] = Field(default_factory=list)


class _VisualizationDemo(BaseModel):
    title: str
    demo_type: str
    widget_type: str = "diagram"
    learning_goal: str
    description: str
    variables: list[str] = Field(min_length=2, max_length=8)
    frames: list[_VisualizationFrame] = Field(min_length=4, max_length=12)
    controls: list[_VisualizationControl]
    teaching_points: list[str] = Field(min_length=3, max_length=8)
    student_tasks: list[str] = Field(min_length=2, max_length=6)
    safety_notes: list[str] = Field(min_length=1)
    nodes: list[dict[str, Any]] = Field(default_factory=list)
    edges: list[dict[str, str]] = Field(default_factory=list)
    code_snippet: str = ""
    physics_scene: Optional[_PhysicsSceneSpec] = None


class _DialogueCard(BaseModel):
    card_type: str
    title: str
    content: str
    metadata: dict[str, Any]


class _DialogueAgentResponse(BaseModel):
    answer: str
    cards: list[_DialogueCard]
    suggested_actions: list[str]
    profile_update_suggestion: str


class _ReflectionEvaluation(BaseModel):
    score: int = Field(ge=0, le=100)
    passed: bool
    feedback: str


class _ResearchReflectionEvaluation(BaseModel):
    detail_score: int = Field(ge=0, le=100)
    relevance_score: int = Field(ge=0, le=100)
    workload_score: int = Field(ge=0, le=100)
    planning_score: int = Field(ge=0, le=100)
    critical_score: int = Field(ge=0, le=100)
    score: int = Field(ge=0, le=100)
    passed: bool
    feedback: str
    strengths: list[str] = Field(default_factory=list)
    improvement_suggestions: list[str] = Field(default_factory=list)
    next_plan_suggestions: list[str] = Field(default_factory=list)


def _item_or_404(db: Session, user: User, item_id: int) -> LearningSyllabusItem:
    item = db.scalar(
        select(LearningSyllabusItem).where(LearningSyllabusItem.id == item_id, LearningSyllabusItem.user_id == user.id)
    )
    if item is None:
        raise KeyError("syllabus item not found")
    return item


def _project_or_404(db: Session, user: User, project_id: int) -> LearningProject:
    project = db.scalar(
        select(LearningProject).where(LearningProject.id == project_id, LearningProject.user_id == user.id)
    )
    if project is None:
        raise KeyError("project not found")
    return project


def _session_query() -> Select[tuple[ClassroomSession]]:
    return select(ClassroomSession).options(
        selectinload(ClassroomSession.resources),
        selectinload(ClassroomSession.submissions),
    )


def _get_session(db: Session, user: User, session_id: int) -> ClassroomSession:
    session = db.scalar(_session_query().where(ClassroomSession.id == session_id, ClassroomSession.user_id == user.id))
    if session is None:
        raise KeyError("classroom session not found")
    return session


def _latest_session(db: Session, user: User, item_id: int) -> Optional[ClassroomSession]:
    return db.scalar(
        _session_query()
        .where(ClassroomSession.syllabus_item_id == item_id, ClassroomSession.user_id == user.id)
        .order_by(ClassroomSession.created_at.desc())
    )


def get_or_create_classroom_session(db: Session, user: User, item_id: int) -> ClassroomSession:
    item = _item_or_404(db, user, item_id)
    project = _project_or_404(db, user, item.project_id)
    session = _latest_session(db, user, item.id)
    if session is None:
        session = ClassroomSession(
            syllabus_item_id=item.id,
            project_id=project.id,
            user_id=user.id,
            title=item.title,
            status="queued",
            progress_state=_build_progress_state(False, False, False, False),
        )
        db.add(session)
        db.add(
            LearningProjectEvent(
                project_id=project.id,
                user_id=user.id,
                event_type="classroom_started",
                summary=f"进入课堂：{item.title}",
                payload={"item_id": item.id},
            )
        )
        db.flush()
    if item.status == "pending":
        update_syllabus_item_status(
            db,
            user,
            item.id,
            SyllabusItemStatusRequest(status="in_progress", reason="进入课堂自动开始学习"),
        )
    db.commit()
    return _get_session(db, user, session.id)


def pre_generate_project_classrooms(db: Session, user: User, project_id: int) -> dict:
    project = _project_or_404(db, user, project_id)
    version = db.scalar(
        select(LearningSyllabusVersion)
        .options(selectinload(LearningSyllabusVersion.items))
        .where(
            LearningSyllabusVersion.project_id == project.id,
            LearningSyllabusVersion.user_id == user.id,
            LearningSyllabusVersion.is_current.is_(True),
            LearningSyllabusVersion.status != "deleted",
        )
        .order_by(LearningSyllabusVersion.version_no.desc())
    )
    if version is None:
        raise KeyError("current syllabus not found")

    items = [
        item
        for item in sorted(version.items, key=lambda value: value.user_order)
        if item.status not in {"deleted", "skipped", "split", "merged"}
    ]
    project.status = "resources_generating"
    project.current_stage = "课堂资源生成中"
    project.next_step = "系统正在按学习清单顺序预生成课堂、课件、例题和复盘资源。"
    _write_event(
        db,
        project,
        user,
        "classroom_resources_generation_started",
        f"开始预生成 {len(items)} 个学习项的课堂资源",
        {"version_id": version.id, "item_count": len(items)},
    )
    db.commit()

    preheat_items = items[:PREGENERATE_CLASSROOM_LIMIT]
    queued_count = max(0, len(items) - len(preheat_items))
    generated = 0
    skipped = 0
    queued = 0
    for item in items:
        session = _get_or_create_prepared_classroom_session(db, user, project, item)
        if item not in preheat_items:
            if not session.ppt_resource_id and session.status not in {"generating", "ready", "completed"}:
                session.status = "queued"
                session.progress_state = _build_progress_state(False, False, False, False)
                db.commit()
            queued += 1
            continue
        try:
            if session.ppt_resource_id:
                skipped += 1
                continue
            generate_classroom_ppt(db, user, session.id, "")
            generated += 1
        except Exception as exc:
            db.rollback()
            project = _project_or_404(db, user, project_id)
            project.status = "resources_failed"
            project.current_stage = "课堂资源生成失败"
            project.next_step = f"{item.title}: {exc.__class__.__name__}: {exc}"
            _write_event(
                db,
                project,
                user,
                "classroom_resources_generation_failed",
                f"课堂资源生成失败：{item.title}",
                {
                    "version_id": version.id,
                    "item_id": item.id,
                    "error": f"{exc.__class__.__name__}: {exc}",
                },
            )
            db.add(
                AgentTaskRecord(
                    session_id=f"project-{project.id}-classroom-resources",
                    user_id=user.id,
                    agent="ClassroomResourcePreGenerationAgent",
                    status="failed",
                    input_summary=f"{project.title} / {item.title}",
                    output_summary=f"{exc.__class__.__name__}: {exc}",
                    latency_ms=0,
                )
            )
            db.commit()
            return {
                "status": "failed",
                "generated": generated,
                "skipped": skipped,
                "queued": queued,
                "failed_item_id": item.id,
                "error": f"{exc.__class__.__name__}: {exc}",
            }

    project = _project_or_404(db, user, project_id)
    project.status = "resources_ready"
    project.current_stage = "核心课堂资源已就绪"
    project.next_step = "已优先生成靠前学习项，后续课堂会在进入学习时继续生成。"
    project.progress = max(project.progress, 15)
    project.generated_resource_count = max(project.generated_resource_count, generated + skipped)
    _write_event(
        db,
        project,
        user,
        "classroom_resources_generated",
        f"课堂资源预热完成：新增 {generated} 个，跳过 {skipped} 个，排队 {queued_count} 个",
        {"version_id": version.id, "generated": generated, "skipped": skipped, "queued": queued},
    )
    db.add(
        AgentTaskRecord(
            session_id=f"project-{project.id}-classroom-resources",
            user_id=user.id,
            agent="ClassroomResourcePreGenerationAgent",
            status="completed",
            input_summary=f"{project.title} / syllabus-v{version.version_no}",
            output_summary=f"预热 {generated} 个课堂资源，跳过 {skipped} 个，排队 {queued_count} 个",
            latency_ms=0,
        )
    )
    db.commit()
    return {"generated": generated, "skipped": skipped, "queued": queued, "total": len(items)}


def _get_or_create_prepared_classroom_session(
    db: Session,
    user: User,
    project: LearningProject,
    item: LearningSyllabusItem,
) -> ClassroomSession:
    session = _latest_session(db, user, item.id)
    if session is not None:
        return session
    session = ClassroomSession(
        syllabus_item_id=item.id,
        project_id=project.id,
        user_id=user.id,
        title=item.title,
        status="queued",
        progress_state=_build_progress_state(False, False, False, False),
    )
    db.add(session)
    db.flush()
    _write_event(
        db,
        project,
        user,
        "classroom_prepared",
        f"预创建课堂：{item.title}",
        {"item_id": item.id, "session_id": session.id},
    )
    db.commit()
    return _get_session(db, user, session.id)


def generate_classroom_ppt(
    db: Session,
    user: User,
    session_id: int,
    instruction: str = "",
) -> ClassroomSession:
    session = _get_session(db, user, session_id)
    item = _item_or_404(db, user, session.syllabus_item_id)
    project = _project_or_404(db, user, session.project_id)
    if session.ppt_resource_id:
        return session
    session.status = "generating"
    session.generation_started_at = datetime.utcnow()
    session.generation_error = ""
    session.progress_state = _build_generation_progress_state("generating", "正在生成课堂课件、例题和复盘任务")
    db.commit()

    package = _generate_classroom_package(project, item, instruction, user=user)
    ppt_path = _write_pptx_file(session.id, item, package)

    resource = ClassroomResource(
        session_id=session.id,
        syllabus_item_id=item.id,
        project_id=project.id,
        user_id=user.id,
        resource_type="pptx",
        title=f"{item.title} 课堂课件",
        content_data=package.model_dump(),
        file_path=str(ppt_path),
        source="THU-MAIC/OpenMAIC-inspired MIT",
        status="ready",
    )
    db.add(resource)
    db.flush()
    session.ppt_resource_id = resource.id
    session.status = "ready"
    session.generation_error = ""
    session.slides_completed = False
    session.slide_progress = {"current_index": 0, "total_slides": len(package.slides), "visited_indices": [0]}
    session.practice_passed = True
    session.progress_state = _build_progress_state(True, False, session.quiz_passed, session.reflection_passed)
    db.add(
        AgentTaskRecord(
            session_id=f"classroom-{session.id}",
            user_id=user.id,
            agent="OpenMAICStyleLessonAgent+PPTAgent",
            status="completed",
            input_summary=f"{project.title} / {item.title}",
            output_summary=f"生成 {len(package.slides)} 页课堂 PPT，并生成例题和复盘要求",
            latency_ms=0,
        )
    )
    _write_event(db, project, user, "classroom_ppt_generated", f"生成课堂 PPT：{item.title}", {"resource_id": resource.id})
    _maybe_complete_session(db, user, session)
    db.commit()
    return _get_session(db, user, session.id)


def request_classroom_ppt_generation(db: Session, user: User, session_id: int) -> tuple[ClassroomSession, bool]:
    session = _get_session(db, user, session_id)
    if session.ppt_resource_id:
        return session, False
    if session.status == "generating":
        return session, False
    session.status = "generating"
    session.generation_started_at = datetime.utcnow()
    session.generation_error = ""
    session.progress_state = _build_generation_progress_state("queued", "课堂资源已进入生成队列")
    db.commit()
    return _get_session(db, user, session.id), True


def mark_classroom_generation_failed(db: Session, user: User, session_id: int, exc: Exception) -> None:
    session = _get_session(db, user, session_id)
    session.status = "failed"
    session.generation_error = f"{exc.__class__.__name__}: {exc}"
    session.progress_state = _build_generation_progress_state("failed", session.generation_error)
    db.commit()


def generate_classroom_voice(
    db: Session,
    user: User,
    session_id: int,
    request: ClassroomVoiceGenerateRequest,
) -> ClassroomSession:
    session = _get_session(db, user, session_id)
    item = _item_or_404(db, user, session.syllabus_item_id)
    project = _project_or_404(db, user, session.project_id)
    package = _classroom_package_or_error(session)
    voice_text = _voice_text_for_scope(package, request)
    voice_path = _write_voice_script_file(session.id, item, request, voice_text)
    resource = ClassroomResource(
        session_id=session.id,
        syllabus_item_id=item.id,
        project_id=project.id,
        user_id=user.id,
        resource_type="voice_script",
        title=f"{item.title} 语音讲解稿",
        content_data={
            "voice_name": request.voice_name,
            "speed": request.speed,
            "text_scope": request.text_scope,
            "slide_index": request.slide_index,
            "page_context": request.page_context,
            "text": voice_text,
            "provider": "browser_speech_synthesis",
            "xunfei_ready": False,
        },
        file_path=str(voice_path),
        source="Browser SpeechSynthesis; Xunfei TTS adapter pending",
        status="ready",
    )
    db.add(resource)
    db.add(
        AgentTaskRecord(
            session_id=f"classroom-{session.id}",
            user_id=user.id,
            agent="VoiceAgent",
            status="completed",
            input_summary=f"{request.text_scope} / {request.voice_name}",
            output_summary=f"生成课堂语音播放稿：{item.title}",
            latency_ms=0,
        )
    )
    _write_event(db, project, user, "classroom_voice_script_generated", f"生成语音讲解稿：{item.title}", {"resource_id": resource.id})
    db.commit()
    return _get_session(db, user, session.id)


def send_classroom_dialogue(
    db: Session,
    user: User,
    session_id: int,
    request: ClassroomDialogueRequest,
) -> dict[str, Any]:
    session = _get_session(db, user, session_id)
    item = _item_or_404(db, user, session.syllabus_item_id)
    project = _project_or_404(db, user, session.project_id)
    package = _classroom_package_or_error(session)
    response = _generate_dialogue_response(project, item, session, package, request, user=user)

    _add_submission(
        db,
        session=session,
        user=user,
        item=item,
        submission_type="dialogue_user",
        content={"role": "user", "message": request.message, "quick_action": request.quick_action},
        score=0,
        passed=True,
        feedback="课堂追问已记录",
    )
    _add_submission(
        db,
        session=session,
        user=user,
        item=item,
        submission_type="dialogue_assistant",
        content={
            "role": "assistant",
            "message": response.answer,
            "cards": [card.model_dump() for card in response.cards],
            "suggested_actions": response.suggested_actions,
            "profile_update_suggestion": response.profile_update_suggestion,
        },
        score=0,
        passed=True,
        feedback="AI 助教已回答",
    )
    db.add(
        AgentTaskRecord(
            session_id=f"classroom-{session.id}",
            user_id=user.id,
            agent="DialogueAgent",
            status="completed",
            input_summary=request.message[:300],
            output_summary=response.answer[:500],
            latency_ms=0,
        )
    )
    db.commit()
    return {
        "answer": response.answer,
        "cards": [card.model_dump() for card in response.cards],
        "suggested_actions": response.suggested_actions,
        "profile_update_suggestion": response.profile_update_suggestion,
        "session": _get_session(db, user, session.id),
    }


def submit_quiz(
    db: Session,
    user: User,
    session_id: int,
    request: ClassroomQuizSubmitRequest,
) -> ClassroomSession:
    session = _get_session(db, user, session_id)
    if not session.slides_completed:
        raise ValueError("请先在页面中翻完课件并完成课件学习")
    item = _item_or_404(db, user, session.syllabus_item_id)
    package = _classroom_package_or_error(session)
    expected = {question.id: question.answer for question in package.quiz}
    results = [_grade_quiz_question(question, request.answers.get(question.id)) for question in package.quiz]
    correct_count = sum(1 for result in results if result["correct"])
    score = round((correct_count / max(len(package.quiz), 1)) * 100)
    passed = score >= 70
    feedback = f"答对 {correct_count}/{len(package.quiz)} 题，{'已达到通过标准' if passed else '请根据提示修改后再次提交'}。"
    _add_submission(
        db,
        session=session,
        user=user,
        item=item,
        submission_type="quiz",
        content={"answers": request.answers, "expected": expected, "results": results, "attempt_no": _next_submission_attempt(session, "quiz")},
        score=score,
        passed=passed,
        feedback=feedback,
    )
    for result in results:
        if result["correct"]:
            continue
        _add_submission(
            db,
            session=session,
            user=user,
            item=item,
            submission_type="mistake",
            content={
                "source": "quiz",
                "question_id": result["question_id"],
                "prompt": result["prompt"],
                "selected": result["selected"],
                "expected": result["expected"],
                "hint": result["hint"],
                "explanation": result["explanation"],
                "knowledge_point": result["knowledge_point"],
            },
            score=0,
            passed=False,
            feedback=result["hint"],
        )
    session.quiz_passed = passed
    session.progress_state = _build_progress_state(bool(session.ppt_resource_id), session.slides_completed, passed, session.reflection_passed)
    _maybe_complete_session(db, user, session)
    db.commit()
    return _get_session(db, user, session.id)


def _grade_quiz_question(question: _QuizSpec, raw_answer: Any) -> dict[str, Any]:
    selected = _normalize_submitted_answer(raw_answer)
    expected = _normalize_submitted_answer(question.answer)
    correct = selected == expected
    return {
        "question_id": question.id,
        "prompt": question.prompt,
        "question_type": question.question_type,
        "options": question.options,
        "selected": ",".join(selected),
        "expected": ",".join(expected),
        "correct": correct,
        "hint": question.hint,
        "explanation": question.explanation,
        "knowledge_point": question.knowledge_point,
    }


def _normalize_submitted_answer(value: Any) -> list[str]:
    parts: list[str] = []
    for item_value in _as_list(value):
        parts.extend(part for part in re.split(r"[,，、\s]+", str(item_value).strip().upper()) if part)
    return sorted(set(parts))


def _next_submission_attempt(session: ClassroomSession, submission_type: str) -> int:
    return 1 + sum(1 for submission in session.submissions if submission.submission_type == submission_type)


def complete_slides(
    db: Session,
    user: User,
    session_id: int,
    request: ClassroomSlidesCompleteRequest,
) -> ClassroomSession:
    session = _get_session(db, user, session_id)
    _classroom_package_or_error(session)
    visited = sorted({index for index in request.visited_indices if 0 <= index < request.total_slides})
    if len(visited) < request.total_slides:
        raise ValueError("课件页尚未全部浏览，不能进入例题环节")
    session.slides_completed = True
    session.slide_progress = {
        "current_index": request.current_index,
        "total_slides": request.total_slides,
        "visited_indices": visited,
        "completed_at": datetime.utcnow().isoformat(),
    }
    session.progress_state = _build_progress_state(True, True, session.quiz_passed, session.reflection_passed)
    item = _item_or_404(db, user, session.syllabus_item_id)
    _write_event(
        db,
        _project_or_404(db, user, session.project_id),
        user,
        "classroom_slides_completed",
        f"完成动态课件学习：{item.title}",
        {"session_id": session.id, "total_slides": request.total_slides},
    )
    db.commit()
    return _get_session(db, user, session.id)


def submit_reflection(
    db: Session,
    user: User,
    session_id: int,
    request: ClassroomReflectionSubmitRequest,
) -> ClassroomSession:
    session = _get_session(db, user, session_id)
    item = _item_or_404(db, user, session.syllabus_item_id)
    project = _project_or_404(db, user, session.project_id)
    if _is_research_reflection(project, item):
        research_evaluation = _evaluate_research_reflection_lenient(project, item, request, user=user)
        evaluation = _ReflectionEvaluation(
            score=research_evaluation.score,
            passed=True,
            feedback=research_evaluation.feedback,
        )
        content = {
            **request.model_dump(),
            "rubric_scores": {
                "详实程度": research_evaluation.detail_score,
                "关联度": research_evaluation.relevance_score,
                "工作量": research_evaluation.workload_score,
                "规划性": research_evaluation.planning_score,
                "批判性思考": research_evaluation.critical_score,
            },
            "strengths": research_evaluation.strengths,
            "improvement_suggestions": research_evaluation.improvement_suggestions,
            "next_plan_suggestions": research_evaluation.next_plan_suggestions,
        }
    else:
        evaluation = _evaluate_reflection_lenient(project, item, request, user=user)
        content = request.model_dump()
    _add_submission(
        db,
        session=session,
        user=user,
        item=item,
        submission_type="reflection",
        content=content,
        score=evaluation.score,
        passed=True,
        feedback=evaluation.feedback,
    )
    session.reflection_passed = True
    session.progress_state = _build_progress_state(bool(session.ppt_resource_id), session.slides_completed, session.quiz_passed, True)
    _maybe_complete_session(db, user, session)
    if session.status == "completed":
        _write_completion_to_profile(db, user, project, item, session, request, evaluation.score)
    db.commit()
    return _get_session(db, user, session.id)


def save_classroom_note(
    db: Session,
    user: User,
    session_id: int,
    request: ClassroomNoteSaveRequest,
) -> ClassroomSession:
    session = _get_session(db, user, session_id)
    item = _item_or_404(db, user, session.syllabus_item_id)
    _classroom_package_or_error(session)
    _add_submission(
        db,
        session=session,
        user=user,
        item=item,
        submission_type="note",
        content=request.model_dump(),
        score=0,
        passed=True,
        feedback="课堂笔记已保存，可用于后续知识库链接与整合。",
    )
    _write_event(
        db,
        _project_or_404(db, user, session.project_id),
        user,
        "classroom_note_saved",
        f"保存课堂笔记：{item.title}",
        {"session_id": session.id, "slide_index": request.slide_index, "slide_title": request.slide_title},
    )
    db.commit()
    return _get_session(db, user, session.id)


def _is_research_reflection(project: LearningProject, item: LearningSyllabusItem) -> bool:
    return bool(project.research_training) or item.item_type in {"paper_reading", "literature_review", "paper_writing", "mock_defense"}


def _evaluate_reflection_lenient(
    project: LearningProject,
    item: LearningSyllabusItem,
    request: ClassroomReflectionSubmitRequest,
    *,
    user: Optional[User] = None,
) -> _ReflectionEvaluation:
    try:
        evaluation = qwen_chat_json(
            "你是高校课程复盘评估助教。请只输出 JSON，给复盘一个演示友好的分数；只要学生写了具体内容，passed 必须为 true。",
            (
                f"项目：{project.title}\n"
                f"学习项：{item.title}\n"
                f"完成标准：{item.completion_criteria}\n"
                f"复盘内容：{request.reflection}\n"
                f"未解决问题：{request.unresolved_questions}\n"
                f"下一步行动：{request.next_action}\n"
                '输出格式：{"score":0-100,"passed":true,"feedback":"..."}'
            ),
            _ReflectionEvaluation,
            user=user,
        )
        score = max(1, int(evaluation.score or 1))
        return _ReflectionEvaluation(score=score, passed=True, feedback=evaluation.feedback or "复盘已记录，课堂学习已完成。")
    except (LLMConfigurationError, LLMResponseError):
        return _ReflectionEvaluation(score=_local_reflection_score(request), passed=True, feedback="复盘已记录，课堂学习已完成。")


def _evaluate_research_reflection_lenient(
    project: LearningProject,
    item: LearningSyllabusItem,
    request: ClassroomReflectionSubmitRequest,
    *,
    user: Optional[User] = None,
) -> _ResearchReflectionEvaluation:
    try:
        evaluation = _evaluate_research_reflection(project, item, request, user=user)
        score = max(1, int(evaluation.score or 1))
        return evaluation.model_copy(update={"score": score, "passed": True})
    except (LLMConfigurationError, LLMResponseError):
        score = _local_reflection_score(request)
        return _ResearchReflectionEvaluation(
            detail_score=score,
            relevance_score=score,
            workload_score=score,
            planning_score=score,
            critical_score=max(1, min(score, 70)),
            score=score,
            passed=True,
            feedback="科研复盘已记录，课堂学习已完成。",
            strengths=["已提交本次阅读/学习复盘"],
            improvement_suggestions=[],
            next_plan_suggestions=[request.next_action] if request.next_action else [],
        )


def _local_reflection_score(request: ClassroomReflectionSubmitRequest) -> int:
    text = " ".join([request.reflection, request.next_action, " ".join(request.unresolved_questions)]).strip()
    return max(1, min(100, 55 + len(text) // 8))


def _evaluate_research_reflection(
    project: LearningProject,
    item: LearningSyllabusItem,
    request: ClassroomReflectionSubmitRequest,
    *,
    user: Optional[User] = None,
) -> _ResearchReflectionEvaluation:
    paper_context = _paper_focus_context(item)
    return qwen_chat_json(
        "你是高校科研训练复盘评估 Agent。请只输出 JSON，并严格按五维量表评分。",
        (
            f"科研项目：{project.title}\n"
            f"科研方向：{project.research_direction}\n"
            f"学习项：{item.title}\n"
            f"单篇论文焦点：{paper_context}\n"
            f"完成标准：{item.completion_criteria}\n"
            f"复盘内容：{request.reflection}\n"
            f"未解决问题：{request.unresolved_questions}\n"
            f"下一步计划：{request.next_action}\n"
            "评分维度：\n"
            "1. detail_score 详实程度：是否具体覆盖论文问题、方法、证据、结论、局限。\n"
            "2. relevance_score 关联度：是否关联学生科研方向、选题边界和后续实验。\n"
            "3. workload_score 工作量：是否体现真实阅读、摘录、对比或复现实操投入。\n"
            "4. planning_score 规划性：下一步计划是否可执行、可检查、有时间/任务边界。\n"
            "5. critical_score 批判性思考：是否指出假设、局限、反例、可改进点或争议。\n"
            "总分 score 为五维加权平均；passed 只有在总分 >= 70 且五维都 >= 50 时为 true。\n"
            "输出字段：detail_score, relevance_score, workload_score, planning_score, critical_score, "
            "score, passed, feedback, strengths, improvement_suggestions, next_plan_suggestions。"
        ),
        _ResearchReflectionEvaluation,
        user=user,
    )


def _write_pptx_file(session_id: int, item: LearningSyllabusItem, package: _ClassroomPackage) -> Path:
    output_dir = Path(__file__).resolve().parents[2] / "generated" / "pptx"
    output_dir.mkdir(parents=True, exist_ok=True)
    path = output_dir / f"classroom-{session_id}-item-{item.id}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank_layout)
    _paint_slide_background(title_slide, prs)
    _add_label(title_slide, "AI CLASSROOM", Inches(0.72), Inches(0.58), Inches(2.2), Inches(0.32), Pt(11), RGBColor(13, 148, 136), bold=True)
    _add_text(
        title_slide,
        package.title,
        Inches(0.72),
        Inches(1.55),
        Inches(8.2),
        Inches(1.4),
        Pt(36),
        RGBColor(19, 78, 74),
        bold=True,
    )
    _add_text(title_slide, item.title, Inches(0.76), Inches(3.05), Inches(8.8), Inches(0.55), Pt(18), RGBColor(71, 85, 105))
    _add_text(title_slide, package.learning_summary, Inches(0.78), Inches(3.85), Inches(6.6), Inches(1.4), Pt(17), RGBColor(15, 118, 110))
    _add_accent_panel(title_slide, Inches(9.25), Inches(1.12), Inches(3.25), Inches(4.95), "个性化课堂", item.knowledge_points[:4] or [item.title])

    for index, spec in enumerate(package.slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _paint_slide_background(slide, prs)
        _add_label(slide, f"{spec.layout.upper()} · {index:02d}", Inches(0.7), Inches(0.46), Inches(2.8), Inches(0.3), Pt(10), RGBColor(13, 148, 136), bold=True)
        _add_text(slide, spec.title, Inches(0.72), Inches(0.82), Inches(8.35), Inches(0.74), Pt(30), RGBColor(19, 78, 74), bold=True)
        _add_visual_concept_panel(slide, spec.visual_hint, Inches(0.78), Inches(1.62), Inches(3.1), Inches(2.1), index)
        _add_visual_block_grid(slide, spec.visual_blocks, Inches(4.18), Inches(1.66), Inches(4.72), Inches(3.0))
        _add_side_panel(slide, spec.side_panel, Inches(9.15), Inches(1.42), Inches(3.18), Inches(3.18))
        _add_takeaway_strip(slide, spec.takeaways, Inches(0.82), Inches(4.05), Inches(8.1), Inches(0.98))
        _add_learning_card(slide, "案例", spec.example, Inches(0.88), Inches(5.2), Inches(2.62), Inches(1.14), RGBColor(240, 253, 250), RGBColor(15, 118, 110))
        _add_learning_card(slide, "易错点", spec.misconception, Inches(3.72), Inches(5.2), Inches(2.62), Inches(1.14), RGBColor(255, 247, 237), RGBColor(194, 65, 12))
        _add_learning_card(slide, "互动", spec.interaction_prompt, Inches(6.56), Inches(5.2), Inches(2.62), Inches(1.14), RGBColor(239, 246, 255), RGBColor(29, 78, 216))
        _add_source_refs(slide, spec.source_refs, Inches(9.42), Inches(4.86), Inches(2.6), Inches(1.28))
        _add_footer(slide, package.title, index, len(package.slides))
        notes = slide.notes_slide.notes_text_frame
        notes.text = spec.speaker_notes

    quiz_slide = prs.slides.add_slide(blank_layout)
    _paint_slide_background(quiz_slide, prs)
    _add_text(quiz_slide, "课堂例题", Inches(0.72), Inches(0.72), Inches(4.8), Inches(0.62), Pt(30), RGBColor(30, 27, 75), bold=True)
    for index, question in enumerate(package.quiz[:4]):
        left = Inches(0.78 + (index % 2) * 6.0)
        top = Inches(1.72 + (index // 2) * 2.15)
        _add_question_card(quiz_slide, f"{question.id}. {question.prompt}", question.explanation, left, top)

    prs.save(path)
    return path


def _voice_text_for_scope(package: _ClassroomPackage, request: ClassroomVoiceGenerateRequest) -> str:
    scope = request.text_scope
    if scope == "all_slides":
        return "\n\n".join([f"{slide.title}。{slide.speaker_notes}" for slide in package.slides])
    if scope == "five_minutes":
        return package.voice_script.five_minutes or package.voice_script.one_minute
    if scope == "one_minute":
        return package.voice_script.one_minute or package.learning_summary
    if not package.slides:
        raise LLMResponseError("课堂包缺少 slides，无法生成当前页语音讲解")
    if request.slide_index >= len(package.slides):
        raise ValueError(f"slide_index 超出范围：{request.slide_index + 1}/{len(package.slides)}")
    slide = package.slides[request.slide_index]
    context = f"\n补充上下文：{request.page_context}" if request.page_context else ""
    return (
        f"现在讲解：{slide.title}。\n"
        f"{slide.speaker_notes}\n"
        f"这个例子可以这样理解：{slide.example}\n"
        f"容易出错的地方是：{slide.misconception}\n"
        f"接下来请思考：{slide.interaction_prompt}\n"
        "这段语音讲解的目标是帮助学生理解当前页，不是逐字朗读幻灯片。"
        f"{context}"
    )


def _write_voice_script_file(
    session_id: int,
    item: LearningSyllabusItem,
    request: ClassroomVoiceGenerateRequest,
    voice_text: str,
) -> Path:
    output_dir = Path(__file__).resolve().parents[2] / "generated" / "voice"
    output_dir.mkdir(parents=True, exist_ok=True)
    path = output_dir / f"classroom-{session_id}-item-{item.id}-{request.text_scope}.txt"
    path.write_text(voice_text, encoding="utf-8")
    return path


def _paint_slide_background(slide: Any, prs: Presentation) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 251, 255)
    bg.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.4), Inches(-0.45), Inches(3.7), Inches(2.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(204, 251, 241)
    accent.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(6.76), Inches(2.1), Inches(0.09))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(13, 148, 136)
    bar.line.fill.background()


def _add_text(slide: Any, text: str, left: Any, top: Any, width: Any, height: Any, size: Any, color: RGBColor, bold: bool = False) -> Any:
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return shape


def _add_label(slide: Any, text: str, left: Any, top: Any, width: Any, height: Any, size: Any, color: RGBColor, bold: bool = False) -> None:
    shape = _add_text(slide, text, left, top, width, height, size, color, bold)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def _add_bullet_card(slide: Any, text: str, left: Any, top: Any, width: Any, height: Any, number: int) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(204, 251, 241)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.12), top + Inches(0.12), Inches(0.34), Inches(0.34))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(13, 148, 136)
    badge.line.fill.background()
    badge_frame = badge.text_frame
    badge_frame.clear()
    badge_paragraph = badge_frame.paragraphs[0]
    badge_paragraph.text = str(number)
    badge_paragraph.alignment = PP_ALIGN.CENTER
    badge_paragraph.font.size = Pt(10)
    badge_paragraph.font.bold = True
    badge_paragraph.font.color.rgb = RGBColor(255, 255, 255)
    _add_text(slide, text, left + Inches(0.6), top + Inches(0.11), width - Inches(0.76), height - Inches(0.12), Pt(14), RGBColor(51, 65, 85))


def _add_visual_concept_panel(slide: Any, visual_hint: str, left: Any, top: Any, width: Any, height: Any, index: int) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(240, 253, 250)
    panel.line.color.rgb = RGBColor(153, 246, 228)
    center_x = left + Inches(1.22)
    center_y = top + Inches(0.34)
    colors = [RGBColor(13, 148, 136), RGBColor(245, 158, 11), RGBColor(37, 99, 235)]
    for node_index, color in enumerate(colors):
        node_left = center_x + Inches((node_index % 2) * 0.78)
        node_top = center_y + Inches(node_index * 0.45)
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, node_left, node_top, Inches(0.52), Inches(0.52))
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
    _add_label(slide, "VISUAL MODEL", left + Inches(0.24), top + Inches(0.22), width - Inches(0.48), Inches(0.28), Pt(9), RGBColor(15, 118, 110), bold=True)
    _add_text(slide, visual_hint, left + Inches(0.24), top + Inches(1.42), width - Inches(0.48), height - Inches(1.58), Pt(12), RGBColor(19, 78, 74))


def _add_visual_block_grid(slide: Any, blocks: list[dict[str, Any]], left: Any, top: Any, width: Any, height: Any) -> None:
    colors = [
        (RGBColor(236, 253, 245), RGBColor(5, 150, 105)),
        (RGBColor(239, 246, 255), RGBColor(37, 99, 235)),
        (RGBColor(255, 247, 237), RGBColor(217, 119, 6)),
        (RGBColor(245, 243, 255), RGBColor(109, 40, 217)),
    ]
    card_width = width / 2 - Inches(0.12)
    card_height = height / 2 - Inches(0.12)
    for index, block in enumerate(blocks[:4]):
        row = index // 2
        col = index % 2
        block_left = left + col * (card_width + Inches(0.24))
        block_top = top + row * (card_height + Inches(0.24))
        fill, accent = colors[index % len(colors)]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, block_left, block_top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        card.line.color.rgb = RGBColor(226, 232, 240)
        _add_label(slide, str(block["type"]).upper(), block_left + Inches(0.18), block_top + Inches(0.12), card_width - Inches(0.36), Inches(0.22), Pt(8), accent, bold=True)
        _add_text(slide, str(block["title"]), block_left + Inches(0.18), block_top + Inches(0.36), card_width - Inches(0.36), Inches(0.32), Pt(13), RGBColor(15, 23, 42), bold=True)
        _add_text(slide, str(block["content"]), block_left + Inches(0.18), block_top + Inches(0.76), card_width - Inches(0.36), card_height - Inches(1.02), Pt(10), RGBColor(51, 65, 85))
        _add_label(slide, str(block["emphasis"])[:42], block_left + Inches(0.18), block_top + card_height - Inches(0.3), card_width - Inches(0.36), Inches(0.18), Pt(8), accent, bold=True)


def _add_side_panel(slide: Any, panel_data: dict[str, Any], left: Any, top: Any, width: Any, height: Any) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(248, 250, 252)
    panel.line.color.rgb = RGBColor(203, 213, 225)
    _add_text(slide, str(panel_data["title"]), left + Inches(0.24), top + Inches(0.22), width - Inches(0.48), Inches(0.36), Pt(15), RGBColor(15, 23, 42), bold=True)
    for index, item in enumerate(panel_data.get("items", [])[:5]):
        y = top + Inches(0.78 + index * 0.46)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.28), y + Inches(0.08), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(13, 148, 136)
        dot.line.fill.background()
        _add_text(slide, str(item), left + Inches(0.5), y, width - Inches(0.72), Inches(0.36), Pt(10), RGBColor(51, 65, 85))


def _add_takeaway_strip(slide: Any, takeaways: list[str], left: Any, top: Any, width: Any, height: Any) -> None:
    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(15, 118, 110)
    strip.line.fill.background()
    _add_label(slide, "TAKEAWAYS", left + Inches(0.24), top + Inches(0.16), Inches(1.4), Inches(0.22), Pt(8), RGBColor(204, 251, 241), bold=True)
    joined = "  ·  ".join(takeaways[:4])
    _add_text(slide, joined, left + Inches(0.24), top + Inches(0.44), width - Inches(0.48), height - Inches(0.5), Pt(12), RGBColor(255, 255, 255), bold=True)


def _add_source_refs(slide: Any, refs: list[Any], left: Any, top: Any, width: Any, height: Any) -> None:
    _add_label(slide, "SOURCES", left, top, width, Inches(0.22), Pt(8), RGBColor(100, 116, 139), bold=True)
    for index, ref in enumerate(refs[:4]):
        if isinstance(ref, dict):
            title = str(ref.get("title", "")).strip()
            url = str(ref.get("url", "")).strip()
            label = f"{index + 1}. {title}" if not url else f"{index + 1}. {title} ({url})"
        else:
            label = f"{index + 1}. {ref}"
        _add_text(slide, label, left, top + Inches(0.3 + index * 0.24), width, Inches(0.2), Pt(8), RGBColor(71, 85, 105))
def _add_speaker_panel(slide: Any, notes: str, interaction_prompt: str = "") -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(1.6), Inches(3.35), Inches(4.55))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(240, 253, 250)
    panel.line.color.rgb = RGBColor(153, 246, 228)
    _add_label(slide, "AI TUTOR TALK", Inches(9.42), Inches(1.86), Inches(2.4), Inches(0.32), Pt(10), RGBColor(15, 118, 110), bold=True)
    _add_text(slide, notes, Inches(9.42), Inches(2.26), Inches(2.75), Inches(2.25), Pt(12), RGBColor(19, 78, 74))
    if interaction_prompt:
        _add_learning_card(
            slide,
            "互动提问",
            interaction_prompt,
            Inches(9.42),
            Inches(4.72),
            Inches(2.75),
            Inches(1.08),
            RGBColor(255, 247, 237),
            RGBColor(194, 65, 12),
        )


def _add_learning_card(slide: Any, title: str, body: str, left: Any, top: Any, width: Any, height: Any, fill: RGBColor, color: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = RGBColor(204, 251, 241)
    _add_text(slide, title, left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.26), Pt(11), color, bold=True)
    _add_text(slide, body, left + Inches(0.18), top + Inches(0.44), width - Inches(0.36), height - Inches(0.5), Pt(11), RGBColor(51, 65, 85))


def _add_accent_panel(slide: Any, left: Any, top: Any, width: Any, height: Any, title: str, items: list[str]) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(238, 242, 255)
    panel.line.color.rgb = RGBColor(199, 210, 254)
    _add_text(slide, title, left + Inches(0.28), top + Inches(0.28), width - Inches(0.5), Inches(0.4), Pt(17), RGBColor(49, 46, 129), bold=True)
    for index, item_text in enumerate(items[:5]):
        _add_text(slide, f"{index + 1}. {item_text}", left + Inches(0.3), top + Inches(0.86 + index * 0.58), width - Inches(0.6), Inches(0.44), Pt(13), RGBColor(71, 85, 105))


def _add_question_card(slide: Any, prompt: str, explanation: str, left: Any, top: Any) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.35), Inches(1.75))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(219, 228, 255)
    _add_text(slide, prompt, left + Inches(0.28), top + Inches(0.22), Inches(4.75), Inches(0.66), Pt(15), RGBColor(30, 27, 75), bold=True)
    _add_text(slide, explanation, left + Inches(0.28), top + Inches(0.95), Inches(4.75), Inches(0.52), Pt(12), RGBColor(71, 85, 105))


def _add_footer(slide: Any, title: str, index: int, total: int) -> None:
    _add_text(slide, title[:48], Inches(0.72), Inches(6.92), Inches(5.2), Inches(0.25), Pt(9), RGBColor(100, 116, 139))
    footer = _add_text(slide, f"{index}/{total}", Inches(12.0), Inches(6.9), Inches(0.55), Inches(0.25), Pt(10), RGBColor(100, 116, 139), bold=True)
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _classroom_package_or_error(session: ClassroomSession) -> _ClassroomPackage:
    ppt_resource = next((resource for resource in session.resources if resource.id == session.ppt_resource_id), None)
    if ppt_resource is None:
        raise ValueError("请先生成课堂 PPT")
    return _ClassroomPackage.model_validate(ppt_resource.content_data)


def _add_submission(
    db: Session,
    *,
    session: ClassroomSession,
    user: User,
    item: LearningSyllabusItem,
    submission_type: str,
    content: dict,
    score: int,
    passed: bool,
    feedback: str,
) -> None:
    db.add(
        ClassroomSubmission(
            session_id=session.id,
            syllabus_item_id=item.id,
            project_id=item.project_id,
            user_id=user.id,
            submission_type=submission_type,
            content=content,
            score=score,
            passed=passed,
            feedback=feedback,
        )
    )


def _build_progress_state(ppt_ready: bool, slides_completed: bool, quiz_passed: bool, reflection_passed: bool) -> dict:
    return {
        "generation_status": "ready" if ppt_ready else "queued",
        "generation_progress": 100 if ppt_ready else 0,
        "generation_message": "课堂资源已就绪" if ppt_ready else "课堂资源等待生成",
        "ppt_ready": ppt_ready,
        "slides_completed": slides_completed,
        "quiz_passed": quiz_passed,
        "practice_passed": True,
        "reflection_passed": reflection_passed,
        "can_complete": all([ppt_ready, slides_completed, quiz_passed, reflection_passed]),
    }


def _build_generation_progress_state(status: str, message: str) -> dict:
    progress_map = {
        "queued": 5,
        "generating": 35,
        "failed": 0,
    }
    return {
        "generation_status": status,
        "generation_progress": progress_map.get(status, 0),
        "generation_message": message,
        "ppt_ready": False,
        "slides_completed": False,
        "quiz_passed": False,
        "practice_passed": True,
        "reflection_passed": False,
        "can_complete": False,
    }


def _maybe_complete_session(db: Session, user: User, session: ClassroomSession) -> None:
    if not all([session.ppt_resource_id, session.slides_completed, session.quiz_passed, session.reflection_passed]):
        session.status = "learning"
        return
    if session.status == "completed":
        return
    session.status = "completed"
    session.completed_at = datetime.utcnow()
    session.practice_passed = True
    session.progress_state = _build_progress_state(True, True, True, True)
    update_syllabus_item_status(
        db,
        user,
        session.syllabus_item_id,
        SyllabusItemStatusRequest(status="completed", reason="课堂课件、例题与复盘均已通过"),
    )


def _write_completion_to_profile(
    db: Session,
    user: User,
    project: LearningProject,
    item: LearningSyllabusItem,
    session: ClassroomSession,
    request: ClassroomReflectionSubmitRequest,
    score: int,
) -> None:
    try:
        from app.services.workspace_service import update_profile_entry

        points = [str(point) for point in (item.knowledge_points or []) if str(point).strip()]
        bullets = [
            f"学习项目：{project.title}",
            f"学习内容：{item.title}",
            f"学习收获：{request.reflection.strip()}",
            f"复盘得分：{score}",
        ]
        if points:
            bullets.append(f"涉及知识点：{'、'.join(points[:8])}")
        if request.unresolved_questions:
            bullets.append(f"待补问题：{'；'.join(request.unresolved_questions[:6])}")
        if request.next_action.strip():
            bullets.append(f"下一步行动：{request.next_action.strip()}")
        if len(request.reflection) < 120:
            bullets.append("学习习惯线索：偏好短复盘和快速推进，适合用清单式下一步任务承接。")
        elif request.unresolved_questions:
            bullets.append("学习习惯线索：会主动暴露疑问，适合追加错因分析和针对性例题。")
        else:
            bullets.append("学习习惯线索：能完成较完整复盘，适合提高任务综合度。")

        update_profile_entry(db, user, ProfileEntryUpdateRequest(
            key=f"learning_record_{session.id}",
            value=bullets,
            source="classroom_completion",
            source_object_id=f"classroom_session:{session.id}",
            confidence=90,
            update_reason=f"完成课堂学习「{item.title}」后写入本次学习记录",
        ))
        if request.unresolved_questions:
            update_profile_entry(db, user, ProfileEntryUpdateRequest(
                key="weak_points",
                value=request.unresolved_questions[:8],
                source="classroom_reflection",
                source_object_id=f"classroom_session:{session.id}",
                confidence=75,
                update_reason="从复盘未解决问题中提取易错点",
            ))
        if request.next_action.strip():
            update_profile_entry(db, user, ProfileEntryUpdateRequest(
                key="interest_direction",
                value=request.next_action.strip(),
                source="classroom_reflection",
                source_object_id=f"classroom_session:{session.id}",
                confidence=70,
                update_reason="从复盘下一步行动中提取学习兴趣方向",
            ))
        if points:
            update_profile_entry(db, user, ProfileEntryUpdateRequest(
                key="mastery",
                value=f"已完成：{item.title}（{'、'.join(points[:5])}）",
                source="classroom_completion",
                source_object_id=f"classroom_session:{session.id}",
                confidence=80,
                update_reason="根据课堂完成情况更新掌握度分布",
            ))
    except Exception:
        return


def _write_event(db: Session, project: LearningProject, user: User, event_type: str, summary: str, payload: dict) -> None:
    db.add(
        LearningProjectEvent(
            project_id=project.id,
            user_id=user.id,
            event_type=event_type,
            summary=summary,
            payload=payload,
        )
    )


def _generate_dialogue_response(
    project: LearningProject,
    item: LearningSyllabusItem,
    session: ClassroomSession,
    package: _ClassroomPackage,
    request: ClassroomDialogueRequest,
    *,
    user: Optional[User] = None,
) -> _DialogueAgentResponse:
    knowledge_context = build_rag_context_for_classroom(project, item, request.message)
    history = [
        submission.content
        for submission in session.submissions
        if submission.submission_type in {"dialogue_user", "dialogue_assistant"}
    ][-10:]
    return qwen_chat_json(
        (
            "你是 DialogueAgent，是课堂中的 AI 助教。"
            "回答必须依托当前课堂上下文，短而可执行；需要时输出 cards 用于前端卡片展示。"
            "只输出 JSON，不输出 Markdown 外壳。"
            f"{FORMULA_OUTPUT_INSTRUCTIONS}"
        ),
        (
            f"项目：{project.title}\n"
            f"研究方向：{project.research_direction}\n"
            f"学习项：{item.title}\n"
            f"学习目标：{item.objective}\n"
            f"课堂摘要：{package.learning_summary}\n"
            f"概念卡：{[card.model_dump() for card in package.concept_cards]}\n"
            f"引导问题：{[question.model_dump() for question in package.guiding_questions]}\n"
            f"课程知识库来源：\n{knowledge_context}\n"
            f"最近对话：{history}\n"
            f"快捷动作：{request.quick_action}\n"
            f"学生问题：{request.message}\n"
            "输出字段：answer, cards, suggested_actions, profile_update_suggestion。"
            "cards 每项包含 card_type, title, content, metadata；card_type 可为 concept, example, code, quiz, diagram, next_step。"
        ),
        _DialogueAgentResponse,
        user=user,
    )


def _normalize_visualization_variable(raw: Any) -> str:
    if isinstance(raw, dict):
        return _require_text(raw.get("label") or raw.get("name") or raw.get("title") or raw.get("id"), "visualization.variables[]")
    return str(raw).strip()


def _normalize_visualization_frame(raw: Any, index: int, variables: list[str]) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError(f"3D 物理演示 JSON frames[{index}] 必须是对象")
    metrics = _numeric_metrics(raw.get("metrics") if isinstance(raw.get("metrics"), dict) else {})
    if not metrics:
        raise LLMResponseError(f"3D 物理演示 JSON frames[{index}] 缺少可渲染的 metrics")
    return {
        "label": _require_text(raw.get("label"), f"visualization.frames[{index}].label"),
        "metrics": metrics,
        "narrative": _require_text(raw.get("narrative"), f"visualization.frames[{index}].narrative"),
    }


def _normalize_visualization_control(raw: Any, index: int) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError(f"3D 物理演示 JSON controls[{index}] 必须是对象")
    return {
        "name": _require_text(raw.get("name"), f"visualization.controls[{index}].name"),
        "label": _require_text(raw.get("label"), f"visualization.controls[{index}].label"),
        "min_value": _require_metric(raw.get("min_value"), f"visualization.controls[{index}].min_value"),
        "max_value": _require_metric(raw.get("max_value"), f"visualization.controls[{index}].max_value"),
        "default_value": _require_metric(raw.get("default_value"), f"visualization.controls[{index}].default_value"),
        "description": _require_text(raw.get("description"), f"visualization.controls[{index}].description"),
    }


def _normalize_demo_type(value: str) -> str:
    allowed = {
        "signal_wave",
        "network_packet",
        "neural_activation",
        "optimization_landscape",
        "sorting_collision",
        "graph_diffusion",
        "physics_system",
        "general_physics",
    }
    normalized = value.strip().lower().replace("-", "_").replace(" ", "_")
    if normalized not in allowed:
        raise LLMResponseError(f"3D 物理演示 JSON demo_type 非法：{value}；允许值：{', '.join(sorted(allowed))}")
    return normalized


def _normalize_physics_scene(raw: Any, demo_type: str) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError("3D 物理演示 JSON 缺少 physics_scene 对象")
    allowed_kinds = {
        "signal_wave",
        "network_packet",
        "neural_activation",
        "optimization_landscape",
        "sorting_collision",
        "graph_diffusion",
        "general_physics",
        "physics_system",
    }
    scene_kind = _require_text(raw.get("scene_kind"), "visualization.physics_scene.scene_kind")
    scene_kind = scene_kind.strip().lower().replace("-", "_").replace(" ", "_")
    if scene_kind not in allowed_kinds:
        raise LLMResponseError(f"3D 物理演示 JSON scene_kind 非法：{scene_kind}；允许值：{', '.join(sorted(allowed_kinds))}")
    camera = raw.get("camera")
    if not isinstance(camera, dict):
        raise LLMResponseError("3D 物理演示 JSON physics_scene.camera 必须是对象")
    objects = [
        _normalize_physics_object(obj, index)
        for index, obj in enumerate(_as_list(raw.get("objects")), start=1)
    ]
    if len(objects) < 4:
        raise LLMResponseError("3D 物理演示 JSON physics_scene.objects 至少需要 4 个对象")
    if len(objects) > 9:
        raise LLMResponseError("3D 物理演示 JSON physics_scene.objects 最多允许 9 个对象")
    annotations_raw = raw.get("annotations")
    annotations: list[dict[str, Any]] = []
    if isinstance(annotations_raw, list):
        for ann_index, ann in enumerate(annotations_raw):
            if not isinstance(ann, dict):
                continue
            try:
                annotations.append({
                    "origin": _vector3(ann.get("origin"), f"visualization.annotations[{ann_index}].origin"),
                    "direction": _vector3(ann.get("direction"), f"visualization.annotations[{ann_index}].direction"),
                    "length": float(ann.get("length", 2.0)),
                    "color": str(ann.get("color", "#fbbf24")),
                })
            except LLMResponseError:
                continue
    return {
        "scene_kind": scene_kind,
        "gravity": _vector3(raw.get("gravity"), "visualization.physics_scene.gravity"),
        "camera": {
            "position": _vector3(camera.get("position"), "visualization.physics_scene.camera.position"),
            "target": _vector3(camera.get("target"), "visualization.physics_scene.camera.target"),
        },
        "objects": objects,
        "annotations": annotations,
    }


def _normalize_physics_object(raw: Any, index: int) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError(f"3D 物理演示 JSON physics_scene.objects[{index}] 必须是对象")
    allowed_shapes = {"sphere", "box", "cylinder", "packet", "node", "torus", "plane"}
    shape = _require_text(raw.get("shape"), f"visualization.physics_scene.objects[{index}].shape")
    shape = shape.strip().lower().replace("-", "_").replace(" ", "_")
    if shape not in allowed_shapes:
        shape = "sphere"  # fallback instead of hard error for LLM-invented shapes
    mass = _require_metric(raw.get("mass"), f"visualization.physics_scene.objects[{index}].mass")
    if mass < 0:
        raise LLMResponseError(f"3D 物理演示 JSON objects[{index}].mass 不能小于 0")
    return {
        "id": _require_text(raw.get("id"), f"visualization.physics_scene.objects[{index}].id"),
        "label": _require_text(raw.get("label"), f"visualization.physics_scene.objects[{index}].label"),
        "role": _require_text(raw.get("role"), f"visualization.physics_scene.objects[{index}].role"),
        "shape": shape,
        "size": _positive_vector3(raw.get("size"), f"visualization.physics_scene.objects[{index}].size"),
        "position": _vector3(raw.get("position"), f"visualization.physics_scene.objects[{index}].position"),
        "velocity": _vector3(raw.get("velocity"), f"visualization.physics_scene.objects[{index}].velocity"),
        "mass": mass,
        "color": _require_color(raw.get("color"), f"visualization.physics_scene.objects[{index}].color"),
        "particle_emitter": bool(raw.get("particle_emitter")),
    }


def _vector3(value: Any, field_path: str) -> list[float]:
    if not isinstance(value, list) or len(value) != 3:
        raise LLMResponseError(f"模型 JSON 字段 {field_path} 必须是长度为 3 的数字数组")
    vector: list[float] = []
    for index, item in enumerate(value):
        if isinstance(item, bool) or not isinstance(item, (int, float)):
            raise LLMResponseError(f"模型 JSON 字段 {field_path}[{index}] 必须是数字")
        vector.append(float(item))
    return vector


def _positive_vector3(value: Any, field_path: str) -> list[float]:
    vector = _vector3(value, field_path)
    if any(item <= 0 for item in vector):
        raise LLMResponseError(f"模型 JSON 字段 {field_path} 必须全部大于 0")
    return vector


def _require_color(value: Any, field_path: str) -> str:
    color = _require_text(value, field_path)
    if not re.fullmatch(r"#[0-9a-fA-F]{6}", color):
        raise LLMResponseError(f"模型 JSON 字段 {field_path} 必须是 #RRGGBB 十六进制颜色")
    return color


def _numeric_metrics(raw: dict[str, Any]) -> dict[str, float]:
    result: dict[str, float] = {}
    for key, value in raw.items():
        metric = _metric_value(value)
        if metric is not None:
            result[str(key)] = metric
    return result


def _metric_value(value: Any) -> Optional[float]:
    if isinstance(value, bool):
        return 1.0 if value else 0.0
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        stripped = value.strip()
        if not stripped:
            return None
        try:
            return float(stripped)
        except ValueError:
            return (sum(ord(ch) for ch in stripped) % 100) / 100
    if isinstance(value, list):
        return min(1.0, len(value) / 10)
    if isinstance(value, dict):
        return min(1.0, len(value) / 10)
    return None


def _require_metric(value: Any, field_path: str) -> float:
    metric = _metric_value(value)
    if metric is None:
        raise LLMResponseError(f"模型 JSON 字段 {field_path} 必须是可解析数值")
    return metric


def _generate_classroom_package(project: LearningProject, item: LearningSyllabusItem, instruction: str, *, user: Optional[User] = None) -> _ClassroomPackage:
    mode_hint = _classroom_mode_hint(project, item)
    paper_focus = _paper_focus_context(item)
    knowledge_context = _truncate_for_prompt(build_rag_context_for_classroom(project, item, instruction), 4500)
    system_prompt = (
        "你是面向高校学生的 OpenMAIC + PPTAgent 风格多智能体课堂生成系统。"
        "只输出严格 JSON，不要 Markdown，不要解释 JSON 之外的内容。"
        "课堂必须由 PlannerAgent、SlideWriterAgent、VisualDesignerAgent、TutorAgent、ExerciseAgent、DemoAgent、SafetyAgent 协同产出。"
        "生成逻辑参考 PPTAgent/DeepPresenter 的思路：先规划每页教学功能，再设计页面布局、视觉层级、素材槽和讲解脚本。"
        "不要输出三行文字式 PPT；每页必须有明确布局、视觉块、侧栏信息和来源引用。"
    )
    user_prompt = (
        "参考 THU-MAIC/OpenMAIC 的课堂产物组织方式和 PPTAgent/DeepPresenter 的课件生成方式："
        "slides、quizzes、PBL/reflection、interactive HTML simulation、presentation planning、visual layout。\n"
        "必须输出顶层字段：title, learning_summary, slides, concept_cards, diagram, guiding_questions, voice_script, "
        "reproduction_demo, readings, quiz, reflection_prompts, safety_notes。\n"
        "禁止只输出 slides、deck、ppt 或 markdown；所有顶层字段必须一次性完整输出。\n"
        "slides: 4 页，每页必须包含 title, layout, bullets, speaker_notes, visual_hint, visual_blocks, side_panel, "
        "takeaways, source_refs, example, misconception, interaction_prompt。"
        "layout 只能使用 cover, split_visual, process_timeline, comparison_matrix, evidence_cards, lab_workbench, defense_panel 中的一种；"
        "bullets 必须是 2-3 条讲得清楚的课堂要点，每条不超过 36 个汉字或 72 个英文字符；speaker_notes 必须是教师口吻讲解，不是照念 bullets；"
        "visual_hint 描述本页适合生成的图解/动画/3D 演示；"
        "visual_blocks 必须是 2 个结构化视觉块，每项必须包含 type, title, content, emphasis，其中 type 只能是 "
        "concept, process, metric, evidence, comparison, formula, code, warning, question；"
        "side_panel 必须包含 title 和 items，items 为 2-3 条短文本；takeaways 为 2-3 条本页结论；"
        "source_refs 必须是 1 个对象，包含 title 和 url；只有确信真实存在的公开链接才填写 url，否则 url 留空并在 title 中写明知识库文档、页码或幻灯片编号；"
        "不得编造 DOI、论文链接、搜索链接、官网链接或空泛来源；"
        "为避免 JSON 截断，所有字符串保持精炼：speaker_notes 不超过 120 字，visual_blocks.content 不超过 60 字，example/misconception/interaction_prompt 不超过 70 字；"
        "example 给出贴近学生项目的具体例子；"
        "misconception 写出常见误区；interaction_prompt 给出本页可以让学生思考或操作的问题。\n"
        "concept_cards: 3 项，每项包含 name, explanation, scenario, misconception, relation_to_project。\n"
        "diagram: 包含 title, diagram_type, mermaid, explanation；mermaid 使用 flowchart TD 或 graph TD。\n"
        "guiding_questions: 3 项，每项包含 prompt, intent, hint。\n"
        "voice_script: 包含 one_minute, five_minutes, segments。\n"
        "reproduction_demo: 包含 title, task, input_format, code_skeleton, steps, expected_output, parameters, common_errors, report_suggestions。\n"
        "readings: 2 项，每项必须是对象并包含 title, why, source, keywords；"
        "source 写课程知识库文档名、上传资料页码、官方文档或真实公开链接，不确定真实链接时写可追溯资料名称，不得编造 DOI 或不存在 URL；"
        "keywords 至少 1 条。\n"
        "quiz: 2-3 项，全部使用选择题；每项包含 id, prompt, question_type, options, answer, explanation, hint, difficulty, knowledge_point。"
        "question_type 只能是 single 或 multiple；options 为 3-5 个对象，每个包含 label 和 text；"
        "single 的 answer 是一个选项 label，例如 A；multiple 的 answer 用英文逗号连接，例如 A,C；"
        "explanation 必须解释为什么正确项正确、为什么常见错误项不对；hint 是学生答错后看的提示，不直接泄露答案。\n"
        "reflection_prompts: 3 条。safety_notes: 至少 1 条。\n"
        "严禁省略字段，严禁只返回 slides。字段缺失会触发一次自动修复；修复仍失败会被系统直接判定为生成失败。\n"
        f"{FORMULA_OUTPUT_INSTRUCTIONS}\n"
        "如果某字段暂时无法确定，必须基于学习项、课程知识库和项目目标生成可执行内容；不得返回空字符串、空数组或占位符。\n"
        "文献综述模式必须给出论文/资料列表、摘要要点、来源字段、对比矩阵和阅读任务；不得编造已发表事实。\n"
        "选题凝练模式必须形成具体题目、研究问题、方法边界、数据与指标、预期贡献和不可做事项。\n"
        "实验助手模式必须生成技术路线、数据采集方案、评价指标、实验变量、图表规范建议和阶段计划。\n"
        "论文写作模式必须围绕课程论文结构、引用规范、图表规范和防幻觉边界设计练习。\n"
        "模拟答辩模式必须在 slides/guiding_questions/quiz/reflection_prompts 中体现开题、中期、答辩问题、追问、评分和修改建议。\n"
        "如果 mode_hint 为 paper_reading，则本次课堂必须是“一篇论文一个 PPT”："
        "slides 围绕同一篇论文组织，依次讲清研究问题、背景脉络、方法/模型、实验证据、局限与启发、对学生选题的下一步计划；"
        "readings 中必须包含当前论文及 1-3 条延伸阅读；quiz/reflection_prompts 必须围绕该论文理解与复盘。\n"
        f"项目：{project.title}\n"
        f"研究方向：{project.research_direction}\n"
        f"学习目标：{project.learning_goal}\n"
        f"学习项：{item.title}\n"
        f"学习项目标：{item.objective}\n"
        f"知识点：{item.knowledge_points}\n"
        f"关联资料：{item.related_documents}\n"
        f"单篇论文焦点：{paper_focus}\n"
        f"课程知识库来源：\n{knowledge_context}\n"
        f"完成标准：{item.completion_criteria}\n"
        f"评估方式：{item.assessment_method}\n"
        f"mode_hint：{mode_hint}\n"
        f"补充要求：{instruction}\n"
    )
    raw = _qwen_chat_raw_json(
        system_prompt,
        user_prompt,
        max_tokens=CLASSROOM_PACKAGE_MAX_TOKENS,
        timeout_seconds=get_settings().qwen_classroom_timeout_seconds,
        user=user,
    )
    try:
        return _validate_classroom_package(raw, project, item)
    except LLMResponseError as first_error:
        repaired = _repair_classroom_package_json(
            system_prompt=system_prompt,
            original_prompt=user_prompt,
            raw=raw,
            validation_error=str(first_error),
            user=user,
        )
        try:
            return _validate_classroom_package(repaired, project, item)
        except LLMResponseError as second_error:
            raise LLMResponseError(
                "课堂包 JSON 自动修复后仍未通过结构校验："
                f"{second_error}；首次错误：{first_error}"
            ) from second_error


def build_rag_context_for_classroom(project: LearningProject, item: LearningSyllabusItem, instruction: str = "") -> str:
    query = "\n".join(
        [
            project.title,
            project.research_direction,
            project.learning_goal,
            item.title,
            item.objective,
            " ".join(item.knowledge_points or []),
            " ".join(item.related_documents or []),
            instruction,
        ]
    )
    db = object_session(project)
    if db is None:
        raise LLMResponseError("课堂生成无法获取数据库会话，不能检索课程知识库")
    return build_rag_context(db, query, limit=4)


def _classroom_mode_hint(project: LearningProject, item: LearningSyllabusItem) -> str:
    text = " ".join(
        [
            project.title,
            project.research_direction,
            project.learning_goal,
            item.title,
            item.item_type,
            " ".join(item.classroom_types or []),
            " ".join(item.knowledge_points or []),
        ]
    ).lower()
    if any(keyword in text for keyword in ["答辩", "defense", "interview", "面试"]):
        return "mock_defense"
    if any(keyword in text for keyword in ["实验助手", "experiment", "数据采集", "技术路线", "变量"]):
        return "experiment_assistant"
    if any(keyword in text for keyword in ["选题", "topic", "研究问题"]):
        return "topic_selection"
    if item.item_type == "paper_reading" or any(keyword in text for keyword in ["paper_ppt", "论文精读", "单篇论文"]):
        return "paper_reading"
    if any(keyword in text for keyword in ["文献", "literature", "综述", "paper"]):
        return "literature_review"
    if any(keyword in text for keyword in ["论文", "写作", "格式", "引用"]):
        return "paper_writing"
    return "general_ai4s_lesson"


def _paper_focus_context(item: LearningSyllabusItem) -> dict[str, str]:
    if item.item_type != "paper_reading":
        return {}
    source = str((item.related_documents or [""])[0])
    title, _, url = source.partition(" | ")
    return {
        "title": title.strip() or item.title,
        "source_url": url.strip(),
        "review_requirement": "学生必须提交论文复盘总结和下一步计划，系统按五维量表评分。",
    }


def _validate_classroom_package(raw: Any, project: LearningProject, item: LearningSyllabusItem) -> _ClassroomPackage:
    normalized = _normalize_classroom_package(raw, project, item)
    try:
        return _ClassroomPackage.model_validate(normalized)
    except Exception as exc:
        raise LLMResponseError(f"课堂包 JSON 归一化后仍未通过结构校验：{exc}") from exc


def _repair_classroom_package_json(
    *,
    system_prompt: str,
    original_prompt: str,
    raw: Any,
    validation_error: str,
    user: Optional[User] = None,
) -> Any:
    repair_prompt = (
        "你刚才输出的课堂包 JSON 未通过系统结构校验。请只返回修复后的完整 JSON 对象，不要 Markdown，不要解释。\n"
        "修复原则：\n"
        "1. 必须保留原课堂主题和原始教学意图，只补齐或改正不合格字段。\n"
        "2. 不得删除必填顶层字段：title, learning_summary, slides, concept_cards, diagram, guiding_questions, "
        "voice_script, reproduction_demo, readings, quiz, reflection_prompts, safety_notes。\n"
        "3. readings 必须是 2-5 个对象，每个对象必须包含 title, why, source, keywords；"
        "source 可以是课程知识库文档名、上传资料页码、官方文档或真实公开链接；不得编造 DOI、论文链接或不存在的网页。\n"
        "4. slides 必须 4 页；每页必须含 title, layout, bullets, speaker_notes, visual_hint, visual_blocks, "
        "side_panel, takeaways, source_refs, example, misconception, interaction_prompt。\n"
        "5. quiz 必须是选择题，至少 2 道；如果原始内容缺少 quiz，请补出最小可用版本；每题必须含 id, prompt, question_type, options, answer, explanation, hint, difficulty, knowledge_point。\n"

        "7. reflection_prompts 至少 3 条；如果原始内容缺少 reflection_prompts，请根据学习项生成 3 条具体复盘问题。\n"
        "8. 所有数组不得为空；不要使用“待补充”“暂无”“N/A”“占位符”。\n"
        f"校验错误：{validation_error}\n"
        f"原始任务要求：{_truncate_for_prompt(original_prompt, 5000)}\n"
        f"待修复 JSON：{_truncate_for_prompt(json.dumps(raw, ensure_ascii=False), 9000)}\n"
    )
    return _qwen_chat_raw_json(
        system_prompt,
        repair_prompt,
        max_tokens=CLASSROOM_REPAIR_MAX_TOKENS,
        temperature=0.1,
        timeout_seconds=get_settings().qwen_classroom_timeout_seconds,
        user=user,
    )


def _qwen_chat_raw_text(
    system_prompt: str,
    user_prompt: str,
    *,
    max_tokens: int = 9000,
    temperature: float = 0.3,
    timeout_seconds: Optional[int] = None,
    user: Optional[User] = None,
) -> str:
    """Call Qwen API without forcing JSON mode — returns raw text (e.g. HTML)."""
    settings = get_settings()
    config = resolve_chat_config(user, timeout_seconds=timeout_seconds or settings.qwen_timeout_seconds)
    validate_qwen_config(user)
    payload = {
        "model": config.model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        "temperature": temperature,
        "max_tokens": max_tokens,
    }
    data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    req = urllib.request.Request(
        f"{config.base_url}/chat/completions",
        data=data,
        headers={"Authorization": f"Bearer {config.api_key}", "Content-Type": "application/json"},
        method="POST",
    )
    try:
        timeout = config.timeout_seconds
        with urllib.request.urlopen(req, timeout=timeout) as response:
            body = json.loads(response.read().decode("utf-8"))
    except (TimeoutError, socket.timeout) as exc:
        raise LLMResponseError(f"{config.provider} 接口请求超时：{timeout} 秒内未返回") from exc
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="ignore")
        raise LLMResponseError(f"{config.provider} 接口返回错误：{exc.code} {detail}") from exc
    except urllib.error.URLError as exc:
        raise LLMResponseError(f"无法连接 {config.provider} 接口：{exc.reason}") from exc
    try:
        return body["choices"][0]["message"]["content"]
    except (KeyError, IndexError, TypeError) as exc:
        raise LLMResponseError(f"{config.provider} 接口响应缺少 choices[0].message.content") from exc


def _qwen_chat_raw_json(
    system_prompt: str,
    user_prompt: str,
    *,
    max_tokens: int = 9000,
    temperature: float = 0.2,
    timeout_seconds: Optional[int] = None,
    user: Optional[User] = None,
) -> Any:
    settings = get_settings()
    config = resolve_chat_config(user, timeout_seconds=timeout_seconds or settings.qwen_timeout_seconds)
    validate_qwen_config(user)
    payload = {
        "model": config.model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        "temperature": temperature,
        "max_tokens": max_tokens,
        "response_format": {"type": "json_object"},
    }
    data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    request = urllib.request.Request(
        f"{config.base_url}/chat/completions",
        data=data,
        headers={"Authorization": f"Bearer {config.api_key}", "Content-Type": "application/json"},
        method="POST",
    )
    try:
        timeout = config.timeout_seconds
        with urllib.request.urlopen(request, timeout=timeout) as response:
            body = json.loads(response.read().decode("utf-8"))
    except (TimeoutError, socket.timeout) as exc:
        raise LLMResponseError(f"{config.provider} 接口请求超时：{timeout} 秒内未返回") from exc
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="ignore")
        raise LLMResponseError(f"{config.provider} 接口返回错误：{exc.code} {detail}") from exc
    except urllib.error.URLError as exc:
        raise LLMResponseError(f"无法连接 {config.provider} 接口：{exc.reason}") from exc
    try:
        content = body["choices"][0]["message"]["content"]
    except (KeyError, IndexError, TypeError) as exc:
        raise LLMResponseError(f"{config.provider} 接口响应缺少 choices[0].message.content") from exc
    return _extract_raw_json(content)


def _extract_raw_json(text: str) -> Any:
    try:
        return parse_llm_json(text)
    except LLMJsonParseError as exc:
        raise LLMResponseError(str(exc)) from exc


def _normalize_classroom_package(raw: Any, project: LearningProject, item: LearningSyllabusItem) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError("课堂包 JSON 顶层必须是对象")
    data = raw
    concept_cards = [_normalize_concept_card(value, item) for value in _as_list(data.get("concept_cards") or data.get("concepts"))]
    if len(concept_cards) < 3:
        raise LLMResponseError("课堂包 JSON 缺少 concept_cards，至少需要 3 张概念卡")
    guiding_questions = [_normalize_guiding_question(value, item) for value in _as_list(data.get("guiding_questions") or data.get("questions_to_think"))]
    if len(guiding_questions) < 3:
        raise LLMResponseError("课堂包 JSON 缺少 guiding_questions，至少需要 3 个引导问题")
    readings_source = (
        data.get("readings")
        or data.get("recommended_readings")
        or data.get("reading_resources")
        or data.get("learning_resources")
        or data.get("resources")
    )
    if isinstance(readings_source, dict):
        readings_source = (
            readings_source.get("readings")
            or readings_source.get("recommended_readings")
            or readings_source.get("items")
            or readings_source.get("resources")
        )
    readings = [_normalize_reading(value, item) for value in _as_list(readings_source)]
    if len(readings) < 2:
        readings.extend(_fallback_readings_for_classroom(data, project, item, 2 - len(readings)))
    if len(readings) < 2:
        raise LLMResponseError("课堂包 JSON 缺少 readings，至少需要 2 条阅读资源")
    slides = [_normalize_slide(slide, index) for index, slide in enumerate(_as_list(data.get("slides") or data.get("ppt") or data.get("deck")), start=1)]
    if len(slides) < 4:
        raise LLMResponseError("课堂包 JSON 缺少 slides，至少需要 4 页")
    slides = slides[:6]

    quiz_raw = data.get("quiz") or data.get("quizzes") or data.get("questions") or data.get("exercises")
    quiz = []
    for index, question in enumerate(_as_list(quiz_raw), start=1):
        try:
            quiz.append(_normalize_quiz(question, index, item))
        except LLMResponseError:
            continue
    if len(quiz) < 2:
        quiz.extend(_fallback_quiz_for_classroom(project, item, 2 - len(quiz)))
    quiz = quiz[:5]

    reflection_prompts = _as_str_list(data.get("reflection_prompts") or data.get("reflection") or data.get("pbl"))
    if len(reflection_prompts) < 3:
        reflection_prompts.extend(_fallback_reflection_prompts_for_classroom(project, item, 3 - len(reflection_prompts)))
    if len(reflection_prompts) < 3:
        raise LLMResponseError("课堂包 JSON 缺少 reflection_prompts，至少需要 3 条")
    safety_notes = _as_str_list(data.get("safety_notes") or data.get("notes"))
    if not safety_notes:
        safety_notes = _fallback_safety_notes_for_classroom(project, item)
    return {
        "title": _require_text(data.get("title"), "classroom.title"),
        "learning_summary": _require_text(data.get("learning_summary") or data.get("summary"), "classroom.learning_summary"),
        "slides": slides,
        "concept_cards": concept_cards[:6],
        "diagram": _normalize_diagram(data.get("diagram") or data.get("mermaid"), item),
        "guiding_questions": guiding_questions[:6],
        "voice_script": _normalize_voice_script(data.get("voice_script") or data.get("script"), item),
        "reproduction_demo": _normalize_reproduction_demo(data.get("reproduction_demo") or data.get("demo"), item),
        "readings": readings[:5],
        "quiz": quiz,
        "reflection_prompts": reflection_prompts[:6],
        "safety_notes": safety_notes,
    }
def _normalize_slide(raw: Any, index: int) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError(f"课堂包 JSON slides[{index}] 必须是对象")
    slide = raw
    visual_blocks = [_normalize_slide_visual_block(value, index, block_index) for block_index, value in enumerate(_as_list(slide.get("visual_blocks")), start=1)]
    if len(visual_blocks) < 2:
        raise LLMResponseError(f"课堂包 JSON slides[{index}].visual_blocks 至少需要 2 个视觉块")
    side_panel = _normalize_slide_side_panel(slide.get("side_panel"), index)
    takeaways = _as_str_list(slide.get("takeaways"))
    if len(takeaways) < 2:
        raise LLMResponseError(f"课堂包 JSON slides[{index}].takeaways 至少需要 2 条")
    bullets = _normalize_slide_bullets(slide, visual_blocks, takeaways, index)
    source_refs = [_normalize_slide_source_ref(value, index, ref_index) for ref_index, value in enumerate(_as_list(slide.get("source_refs") or slide.get("sources")), start=1)]
    if not source_refs:
        raise LLMResponseError(f"课堂包 JSON slides[{index}].source_refs 至少需要 1 条")
    return {
        "title": _require_text(slide.get("title"), f"classroom.slides[{index}].title"),
        "layout": _normalize_slide_layout(slide.get("layout"), index),
        "bullets": bullets[:6],
        "speaker_notes": _require_text(slide.get("speaker_notes") or slide.get("notes"), f"classroom.slides[{index}].speaker_notes"),
        "visual_hint": _require_text(slide.get("visual_hint") or slide.get("visual") or slide.get("diagram_hint"), f"classroom.slides[{index}].visual_hint"),
        "visual_blocks": visual_blocks[:6],
        "side_panel": side_panel,
        "takeaways": takeaways[:5],
        "source_refs": source_refs[:5],
        "example": _require_text(slide.get("example") or slide.get("case"), f"classroom.slides[{index}].example"),
        "misconception": _require_text(slide.get("misconception") or slide.get("pitfall"), f"classroom.slides[{index}].misconception"),
        "interaction_prompt": _require_text(slide.get("interaction_prompt") or slide.get("question") or slide.get("prompt"), f"classroom.slides[{index}].interaction_prompt"),
    }


def _normalize_slide_layout(value: Any, index: int) -> str:
    layout = _require_text(value, f"classroom.slides[{index}].layout")
    allowed = {"cover", "split_visual", "process_timeline", "comparison_matrix", "evidence_cards", "lab_workbench", "defense_panel"}
    if layout not in allowed:
        raise LLMResponseError(f"classroom.slides[{index}].layout 必须是 {sorted(allowed)} 之一，当前为 {layout}")
    return layout


def _normalize_slide_bullets(slide: dict[str, Any], visual_blocks: list[dict[str, str]], takeaways: list[str], index: int) -> list[str]:
    bullets = _as_str_list(slide.get("bullets") or slide.get("points") or slide.get("content") or slide.get("items"))
    if len(bullets) < 2:
        bullets.extend(str(block.get("content", "")).strip() for block in visual_blocks)
    if len([item for item in bullets if item.strip()]) < 2:
        bullets.extend(takeaways)
    if len([item for item in bullets if item.strip()]) < 2:
        notes = str(slide.get("speaker_notes") or slide.get("notes") or "").strip()
        if notes:
            bullets.extend(_split_compact_sentences(notes))
    cleaned: list[str] = []
    for bullet in bullets:
        text = _compact_text(str(bullet), 110)
        if text and text not in cleaned:
            cleaned.append(text)
    if len(cleaned) < 2:
        raise LLMResponseError(f"课堂包 JSON slides[{index}] 缺少可用于课堂展示的要点；请让模型为 bullets、visual_blocks 或 takeaways 提供至少 2 条具体内容")
    return cleaned[:6]


def _normalize_slide_visual_block(raw: Any, slide_index: int, block_index: int) -> dict[str, str]:
    if not isinstance(raw, dict):
        raise LLMResponseError(f"classroom.slides[{slide_index}].visual_blocks[{block_index}] 必须是对象")
    block_type = _require_text(raw.get("type"), f"classroom.slides[{slide_index}].visual_blocks[{block_index}].type")
    allowed = {"concept", "process", "metric", "evidence", "comparison", "formula", "code", "warning", "question"}
    if block_type not in allowed:
        raise LLMResponseError(
            f"classroom.slides[{slide_index}].visual_blocks[{block_index}].type 必须是 {sorted(allowed)} 之一，当前为 {block_type}"
        )
    return {
        "type": block_type,
        "title": _require_text(raw.get("title"), f"classroom.slides[{slide_index}].visual_blocks[{block_index}].title"),
        "content": _require_text(raw.get("content"), f"classroom.slides[{slide_index}].visual_blocks[{block_index}].content"),
        "emphasis": _require_text(raw.get("emphasis"), f"classroom.slides[{slide_index}].visual_blocks[{block_index}].emphasis"),
    }


def _normalize_slide_side_panel(raw: Any, slide_index: int) -> dict[str, Any]:
    if not isinstance(raw, dict):
        raise LLMResponseError(f"classroom.slides[{slide_index}].side_panel 必须是对象")
    items = _as_str_list(raw.get("items"))
    if len(items) < 2:
        raise LLMResponseError(f"classroom.slides[{slide_index}].side_panel.items 至少需要 2 条")
    return {
        "title": _require_text(raw.get("title"), f"classroom.slides[{slide_index}].side_panel.title"),
        "items": items[:5],
    }


def _normalize_slide_source_ref(raw: Any, slide_index: int, ref_index: int) -> dict[str, str]:
    if isinstance(raw, dict):
        title = _require_text(raw.get("title") or raw.get("name"), f"classroom.slides[{slide_index}].source_refs[{ref_index}].title")
        url = str(raw.get("url") or raw.get("source") or raw.get("uri") or "").strip()
        if url and not re.match(r"^https?://", url, flags=re.IGNORECASE):
            url = ""
        return {"title": title, "url": url}
    text = _compact_text(str(raw), 120)
    if not text:
        raise LLMResponseError(f"classroom.slides[{slide_index}].source_refs[{ref_index}] 必须是对象")
    return {"title": text, "url": ""}
def _normalize_quiz(raw: Any, index: int, item: LearningSyllabusItem) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError(f"课堂包 JSON quiz[{index}] 必须是对象")
    question = raw
    options = _normalize_quiz_options(question.get("options") or question.get("choices"), index)
    option_labels = {option["label"] for option in options}
    question_type = str(question.get("question_type") or question.get("type") or "single").strip().lower()
    if question_type in {"single_choice", "choice", "radio"}:
        question_type = "single"
    if question_type in {"multiple_choice", "multi", "checkbox"}:
        question_type = "multiple"
    if question_type not in {"single", "multiple"}:
        raise LLMResponseError(f"classroom.quiz[{index}].question_type 必须是 single 或 multiple")
    answer = _normalize_quiz_answer(question.get("answer") or question.get("expected_answer") or question.get("correct_answer"), option_labels, question_type, index)
    return {
        "id": _require_text(question.get("id"), f"classroom.quiz[{index}].id"),
        "prompt": _require_text(question.get("prompt") or question.get("question"), f"classroom.quiz[{index}].prompt"),
        "question_type": question_type,
        "options": options,
        "answer": answer,
        "explanation": _require_text(question.get("explanation") or question.get("analysis"), f"classroom.quiz[{index}].explanation"),
        "hint": _require_text(question.get("hint") or question.get("tip") or question.get("feedback_hint"), f"classroom.quiz[{index}].hint"),
        "difficulty": str(question.get("difficulty") or "medium").strip()[:24],
        "knowledge_point": str(question.get("knowledge_point") or question.get("point") or (item.knowledge_points[0] if item.knowledge_points else item.title)).strip()[:80],
    }


def _normalize_quiz_options(raw: Any, index: int) -> list[dict[str, str]]:
    values = _as_list(raw)
    if len(values) < 3:
        raise LLMResponseError(f"classroom.quiz[{index}].options 至少需要 3 个选项")
    labels = "ABCDEFGHIJKLMNOPQRSTUVWXYZ"
    options: list[dict[str, str]] = []
    for option_index, value in enumerate(values[:6]):
        if isinstance(value, dict):
            label = str(value.get("label") or value.get("key") or value.get("id") or labels[option_index]).strip().upper()
            text = _require_text(value.get("text") or value.get("content") or value.get("label_text") or value.get("value"), f"classroom.quiz[{index}].options[{option_index + 1}].text")
        else:
            label = labels[option_index]
            text = str(value).strip()
        if not re.fullmatch(r"[A-Z]", label):
            raise LLMResponseError(f"classroom.quiz[{index}].options[{option_index + 1}].label 必须是 A-Z 单个字母")
        if not text:
            raise LLMResponseError(f"classroom.quiz[{index}].options[{option_index + 1}].text 不能为空")
        options.append({"label": label, "text": _compact_text(text, 140)})
    if len({option["label"] for option in options}) != len(options):
        raise LLMResponseError(f"classroom.quiz[{index}].options 存在重复 label")
    return options


def _normalize_quiz_answer(raw: Any, option_labels: set[str], question_type: str, index: int) -> str:
    answers: list[str] = []
    for value in _as_list(raw):
        if isinstance(value, str):
            parts = re.split(r"[,，、\s]+", value.strip())
            answers.extend(part for part in parts if part)
        else:
            answers.append(str(value).strip())
    normalized = [answer.strip().upper() for answer in answers if answer.strip()]
    if not normalized:
        raise LLMResponseError(f"classroom.quiz[{index}].answer 不能为空")
    invalid = [answer for answer in normalized if answer not in option_labels]
    if invalid:
        raise LLMResponseError(f"classroom.quiz[{index}].answer 引用了不存在的选项：{', '.join(invalid)}")
    if question_type == "single" and len(normalized) != 1:
        raise LLMResponseError(f"classroom.quiz[{index}] 是单选题，answer 只能包含 1 个选项")
    return ",".join(sorted(set(normalized)))


def _normalize_concept_card(raw: Any, item: LearningSyllabusItem) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError("课堂包 JSON concept_cards[] 必须是对象")
    value = raw
    name = _require_text(value.get("name") or value.get("title"), "classroom.concept_cards[].name")
    return {
        "name": name,
        "explanation": _require_text(value.get("explanation") or value.get("description"), "classroom.concept_cards[].explanation"),
        "scenario": _require_text(value.get("scenario") or value.get("use_case"), "classroom.concept_cards[].scenario"),
        "misconception": _require_text(value.get("misconception") or value.get("pitfall"), "classroom.concept_cards[].misconception"),
        "relation_to_project": _require_text(value.get("relation_to_project") or value.get("project_relation"), "classroom.concept_cards[].relation_to_project"),
    }


def _normalize_diagram(raw: Any, item: LearningSyllabusItem) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError("课堂包 JSON diagram 必须是对象")
    value = raw
    mermaid = _require_text(value.get("mermaid") or value.get("source"), "classroom.diagram.mermaid")
    return {
        "title": _require_text(value.get("title"), "classroom.diagram.title"),
        "diagram_type": _require_text(value.get("diagram_type") or value.get("type"), "classroom.diagram.diagram_type"),
        "mermaid": mermaid,
        "explanation": _require_text(value.get("explanation"), "classroom.diagram.explanation"),
    }


def _normalize_voice_script(raw: Any, item: LearningSyllabusItem) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError("课堂包 JSON voice_script 必须是对象")
    value = raw
    segments = _as_str_list(value.get("segments"))
    if not segments:
        raise LLMResponseError("课堂包 JSON voice_script.segments 至少需要 1 条")
    return {
        "one_minute": _require_text(value.get("one_minute"), "classroom.voice_script.one_minute"),
        "five_minutes": _require_text(value.get("five_minutes"), "classroom.voice_script.five_minutes"),
        "segments": segments[:8],
    }


def _normalize_reproduction_demo(raw: Any, item: LearningSyllabusItem) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError("课堂包 JSON reproduction_demo 必须是对象")
    value = raw
    steps = _as_str_list(value.get("steps"))
    if not steps:
        raise LLMResponseError("课堂包 JSON reproduction_demo.steps 至少需要 1 条")
    parameters = _as_str_list(value.get("parameters"))
    if not parameters:
        raise LLMResponseError("课堂包 JSON reproduction_demo.parameters 至少需要 1 条")
    common_errors = _as_str_list(value.get("common_errors") or value.get("errors"))
    if not common_errors:
        raise LLMResponseError("课堂包 JSON reproduction_demo.common_errors 至少需要 1 条")
    report_suggestions = _as_str_list(value.get("report_suggestions") or value.get("report"))
    if not report_suggestions:
        raise LLMResponseError("课堂包 JSON reproduction_demo.report_suggestions 至少需要 1 条")
    return {
        "title": _require_text(value.get("title") or value.get("name"), "classroom.reproduction_demo.title"),
        "task": _require_text(value.get("task"), "classroom.reproduction_demo.task"),
        "input_format": _require_text(value.get("input_format"), "classroom.reproduction_demo.input_format"),
        "code_skeleton": str(value.get("code_skeleton") or value.get("code") or ""),
        "steps": steps[:8],
        "expected_output": _require_text(value.get("expected_output"), "classroom.reproduction_demo.expected_output"),
        "parameters": parameters[:8],
        "common_errors": common_errors[:8],
        "report_suggestions": report_suggestions[:8],
    }


def _normalize_guiding_question(raw: Any, item: LearningSyllabusItem) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError("课堂包 JSON guiding_questions[] 必须是对象")
    value = raw
    return {
        "prompt": _require_text(value.get("prompt") or value.get("question"), "classroom.guiding_questions[].prompt"),
        "intent": _require_text(value.get("intent"), "classroom.guiding_questions[].intent"),
        "hint": _require_text(value.get("hint"), "classroom.guiding_questions[].hint"),
    }


def _normalize_reading(raw: Any, item: LearningSyllabusItem) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError("课堂包 JSON readings[] 必须是对象")
    value = raw
    keywords = _as_str_list(value.get("keywords"))
    if not keywords:
        raise LLMResponseError("课堂包 JSON readings[].keywords 至少需要 1 条")
    return {
        "title": _require_text(value.get("title"), "classroom.readings[].title"),
        "why": _require_text(value.get("why") or value.get("reason"), "classroom.readings[].why"),
        "source": _require_text(value.get("source") or value.get("uri"), "classroom.readings[].source"),
        "keywords": keywords[:8],
    }


def _collect_slide_source_refs(data: dict[str, Any]) -> list[dict[str, str]]:
    refs: list[dict[str, str]] = []
    for slide in _as_list(data.get("slides") or data.get("ppt") or data.get("deck")):
        if not isinstance(slide, dict):
            continue
        for raw in _as_list(slide.get("source_refs") or slide.get("sources")):
            if isinstance(raw, dict):
                title = _compact_text(str(raw.get("title") or raw.get("name") or ""), 90)
                url = str(raw.get("url") or raw.get("source") or raw.get("uri") or "").strip()
            else:
                title = _compact_text(str(raw), 90)
                url = ""
            if not title and not url:
                continue
            refs.append({"title": title or url, "source": title or url, "url": url})
    return refs


def _fallback_keywords(project: LearningProject, item: LearningSyllabusItem, title: str) -> list[str]:
    values = [
        item.title,
        item.item_type,
        project.title,
        project.research_direction,
        project.learning_goal,
        title,
        *list(item.knowledge_points or []),
        *list(item.related_documents or []),
    ]
    keywords: list[str] = []
    for value in values:
        for token in re.split(r"[、,，/|;；\s]+", str(value or "")):
            text = _compact_text(token, 30)
            if text and text not in keywords:
                keywords.append(text)
    return (keywords or [_compact_text(item.title, 20) or "课堂导读"])[:8]


def _fallback_readings_for_classroom(data: dict[str, Any], project: LearningProject, item: LearningSyllabusItem, needed: int) -> list[dict[str, str]]:
    candidates: list[dict[str, str]] = []
    seen: set[str] = set()
    for ref in _collect_slide_source_refs(data):
        key = (ref.get("title") or ref.get("source") or "").strip().lower()
        if not key or key in seen:
            continue
        seen.add(key)
        title = ref.get("title") or item.title
        source = ref.get("url") or ref.get("source") or title
        candidates.append(
            {
                "title": _compact_text(title, 90),
                "why": _compact_text(f"用于补齐课堂包阅读资源，帮助学生先理解 {item.title} 的课件来源。", 180),
                "source": _compact_text(source, 180),
                "keywords": _fallback_keywords(project, item, title),
            }
        )
        if len(candidates) >= needed:
            return candidates

    for title in [f"{item.title} 导读", f"{project.title or item.title} 延伸阅读", project.research_direction]:
        text = _compact_text(str(title or ""), 90)
        key = text.lower()
        if not text or key in seen:
            continue
        seen.add(key)
        candidates.append(
            {
                "title": text,
                "why": _compact_text(f"用于为 {item.title} 补齐最小阅读资源并支撑课堂讲解。", 180),
                "source": _compact_text(project.title or item.title or "课堂知识库", 180),
                "keywords": _fallback_keywords(project, item, text),
            }
        )
        if len(candidates) >= needed:
            break
    return candidates


def _fallback_quiz_for_classroom(project: LearningProject, item: LearningSyllabusItem, needed: int) -> list[dict[str, Any]]:
    points = [str(value).strip() for value in (item.knowledge_points or []) if str(value).strip()]
    focus = points[0] if points else item.title
    templates = [
        {
            "prompt": f"学习 {item.title} 时，最应该先确认哪一项？",
            "options": [
                {"label": "A", "text": f"核心概念：{focus}"},
                {"label": "B", "text": "随机记忆所有材料原文"},
                {"label": "C", "text": "跳过前置条件直接做综合项目"},
            ],
            "answer": "A",
            "explanation": f"课堂包缺少完整测验时，系统用学习项知识点生成最小检测题；本题用于确认学生是否抓住 {focus}。",
            "hint": "先找本节课最核心的知识点。",
            "knowledge_point": focus,
        },
        {
            "prompt": f"完成 {item.title} 后，哪种产出最能证明你已经理解？",
            "options": [
                {"label": "A", "text": "能用自己的话说明概念、步骤和一个例子"},
                {"label": "B", "text": "只保存课件截图"},
                {"label": "C", "text": "只记住课程标题"},
            ],
            "answer": "A",
            "explanation": "理解需要可复述、可举例、可操作的证据，而不是只保留材料表面信息。",
            "hint": "看哪一项能形成可检查的学习证据。",
            "knowledge_point": focus,
        },
    ]
    result: list[dict[str, Any]] = []
    for index, template in enumerate(templates[: max(0, needed)], start=1):
        result.append(
            {
                "id": f"Q{index}",
                "prompt": template["prompt"],
                "question_type": "single",
                "options": template["options"],
                "answer": template["answer"],
                "explanation": template["explanation"],
                "hint": template["hint"],
                "difficulty": "easy",
                "knowledge_point": template["knowledge_point"],
            }
        )
    return result


def _fallback_reflection_prompts_for_classroom(project: LearningProject, item: LearningSyllabusItem, needed: int) -> list[str]:
    focus = [str(value).strip() for value in (item.knowledge_points or []) if str(value).strip()]
    lead = focus[0] if focus else item.title
    prompts = [
        f"今天学到的 {lead} 用一句话怎么解释给同学听？",
        f"这个知识点最容易和什么概念混淆，你会怎么区分？",
        f"如果要继续学习 {project.research_direction or project.title}，下一步你最想补哪一块？",
    ]
    return prompts[: max(0, needed)]


def _fallback_safety_notes_for_classroom(project: LearningProject, item: LearningSyllabusItem) -> list[str]:
    focus = [str(value).strip() for value in (item.knowledge_points or []) if str(value).strip()]
    lead = focus[0] if focus else item.title
    return [
        f"先确认 {lead} 的前置条件，再做练习。",
        "不要把课件里的示例直接当成全部结论。",
        "遇到不确定的地方，先回看知识库来源再提交。",
    ]


def _truncate_for_prompt(value: str, max_length: int) -> str:
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]+", "", value)
    if len(text) <= max_length:
        return text
    return text[:max_length] + "\n...[已截断，保留前文供修复 JSON 结构使用]"


def _as_list(value: Any) -> list[Any]:
    if value is None:
        return []
    if isinstance(value, list):
        return value
    return [value]


def _as_str_list(value: Any) -> list[str]:
    result: list[str] = []
    for item_value in _as_list(value):
        if isinstance(item_value, dict):
            text = item_value.get("text") or item_value.get("content") or item_value.get("title") or json.dumps(item_value, ensure_ascii=False)
        else:
            text = str(item_value)
        text = text.strip()
        if text:
            result.append(text)
    return result


def _compact_text(value: str, max_length: int) -> str:
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]+", "", value)
    text = re.sub(r"\s+", " ", text).strip()
    return text[:max_length].rstrip()


def _split_compact_sentences(value: str) -> list[str]:
    parts = re.split(r"[。！？!?；;]\s*", value)
    return [_compact_text(part, 110) for part in parts if _compact_text(part, 110)]


def _require_text(value: Any, field_path: str) -> str:
    texts = _as_str_list(value)
    if not texts:
        raise LLMResponseError(f"模型 JSON 缺少必填字段 {field_path}")
    return _compact_text(texts[0], 500)


# OpenMAIC-style interactive widget generation.
# Keep this section after the legacy helpers so Python binds these names last.

_VISUAL_WIDGET_TYPES = {"diagram", "simulation", "code", "visualization3d"}
_NON_3D_DEMO_TYPES = {
    "diagram": {"concept_map", "system_diagram", "flowchart", "comparison_map"},
    "simulation": {"state_machine", "data_flow", "algorithm_trace", "process_simulation"},
    "code": {"code_walkthrough", "debug_trace", "api_flow", "reproduction_demo"},
}
_VISUAL_COLOR_PALETTE = ["#4f8b78", "#6f7fb7", "#d08a4f", "#8d6ab8", "#2f7ea8", "#b75f5f", "#6d8f46", "#c28b2e"]
_DEMO_TYPE_ALIASES = {
    "diagram": {
        "concept_graph": "concept_map",
        "knowledge_graph": "concept_map",
        "mind_map": "concept_map",
        "mindmap": "concept_map",
        "architecture": "system_diagram",
        "system_map": "system_diagram",
        "flow": "flowchart",
        "flow_chart": "flowchart",
        "comparison": "comparison_map",
    },
    "simulation": {
        "algorithm_execution_trace": "algorithm_trace",
        "execution_trace": "algorithm_trace",
        "algorithm_execution": "algorithm_trace",
        "algorithm_simulation": "algorithm_trace",
        "algorithm_flow": "algorithm_trace",
        "step_trace": "algorithm_trace",
        "trace": "algorithm_trace",
        "dataflow": "data_flow",
        "flow": "data_flow",
        "data_pipeline": "data_flow",
        "process": "process_simulation",
        "process_flow": "process_simulation",
        "fsm": "state_machine",
        "finite_state_machine": "state_machine",
    },
    "code": {
        "code_trace": "code_walkthrough",
        "walkthrough": "code_walkthrough",
        "debugging": "debug_trace",
        "debug": "debug_trace",
        "api": "api_flow",
        "reproduce": "reproduction_demo",
        "repro_demo": "reproduction_demo",
    },
}

# When user forces widget_type → visualization3d, map the LLM's non-3D demo_type to a sensible 3D fallback
_DEMO_TO_3D_FALLBACK = {
    "concept_map": "graph_diffusion",
    "system_diagram": "physics_system",
    "flowchart": "network_packet",
    "comparison_map": "optimization_landscape",
    "algorithm_trace": "sorting_collision",
    "data_flow": "network_packet",
    "process_simulation": "physics_system",
    "state_machine": "sorting_collision",
    "code_walkthrough": "signal_wave",
    "debug_trace": "signal_wave",
    "api_flow": "network_packet",
    "reproduction_demo": "signal_wave",
}

# When user forces widget_type away from visualization3d, map 3D demo_types
_3D_TO_NON_3D_FALLBACK = {
    "diagram": {
        "signal_wave": "concept_map",
        "network_packet": "system_diagram",
        "neural_activation": "concept_map",
        "optimization_landscape": "comparison_map",
        "sorting_collision": "flowchart",
        "graph_diffusion": "concept_map",
        "physics_system": "system_diagram",
    },
    "simulation": {
        "signal_wave": "algorithm_trace",
        "network_packet": "data_flow",
        "neural_activation": "process_simulation",
        "optimization_landscape": "algorithm_trace",
        "sorting_collision": "algorithm_trace",
        "graph_diffusion": "data_flow",
        "physics_system": "process_simulation",
    },
    "code": {
        "signal_wave": "code_walkthrough",
        "network_packet": "api_flow",
        "neural_activation": "code_walkthrough",
        "optimization_landscape": "debug_trace",
        "sorting_collision": "code_walkthrough",
        "graph_diffusion": "api_flow",
        "physics_system": "reproduction_demo",
    },
}


def _generate_visualization3d_html_direct(
    project: LearningProject,
    item: LearningSyllabusItem,
    package: _ClassroomPackage,
    instruction: str,
    *,
    user: Optional[User] = None,
) -> tuple[str, str]:
    """OpenMAIC-style: LLM directly outputs a complete self-contained Three.js HTML page."""
    knowledge_context = build_rag_context_for_classroom(project, item, instruction)
    system_prompt = (
        "你是 3D 可视化专家。根据课堂内容生成一个完整的、自包含的 HTML 文件，使用 Three.js 展示交互式 3D 场景。"
        "只输出完整 HTML 文档（从 <!DOCTYPE html> 到 </html>），不输出 Markdown、解释或代码块标记。"
        f"{FORMULA_OUTPUT_INSTRUCTIONS}"
    )
    user_prompt = (
        '生成一个完整的交互式 3D 演示 HTML 页面。\n\n'
        '## 技术要求\n'
        '1. 使用 Three.js ES 模块 (importmap 从 esm.sh 加载 three@0.167)\n'
        '2. 必须包含 OrbitControls (鼠标旋转/缩放/平移)\n'
        '3. 良好的光照: AmbientLight + DirectionalLight + 可选的 HemisphereLight\n'
        '4. 响应式设计，铺满整个窗口\n'
        '5. 底部控制栏: 播放/暂停按钮、速度滑块、重置按钮、步骤跳转按钮\n'
        '6. 左上角信息面板: 标题和当前步骤说明\n'
        '7. 深色背景 (#0b1120)，对象颜色鲜明可辨\n\n'
        '## 场景要求\n'
        f'根据以下内容设计 3D 场景:\n'
        f'- 项目: {project.title}\n'
        f'- 学习项: {item.title}\n'
        f'- 知识点: {item.knowledge_points}\n'
        f'- 用户要求: {instruction}\n'
        f'- 课堂摘要: {package.learning_summary}\n\n'
        '## 场景设计指南\n'
        '- 创建 4-8 个 3D 对象 (球体/立方体/圆柱/环)，用颜色区分角色\n'
        '- 对象之间用线条连接表示关系\n'
        '- 用 CSS positioned div 做标签，通过 project+transform 定位\n'
        '- 设计 4-6 个步骤帧，存储在 frames 数组中，每帧有 title/description/highlight 字段\n'
        '- 对象在各帧之间可以有位置/颜色/大小变化\n'
        '- 至少有一个对象带有粒子效果 (用 Points 实现)\n\n'
        '## 交互要求\n'
        '- 播放/暂停按钮自动循环步骤 (每步约3秒)\n'
        '- 速度滑块控制播放速度\n'
        '- 鼠标悬停 3D 对象时高亮并显示 tooltip\n'
        '- 底部步骤按钮可直接跳转到指定步骤\n\n'
        '## 关键代码规范 (必须遵守，否则页面白屏)\n'
        '1. 用静态 import，不要用动态 import() -- 动态导入在 importmap 环境下可能失败:\n'
        '   import * as THREE from “three”;\n'
        '   import { OrbitControls } from “three/addons/controls/OrbitControls.js”;\n'
        '2. 渲染循环必须始终运行 -- animate() 无条件调用 requestAnimationFrame 和 renderer.render():\n'
        '   function animate() {\n'
        '       requestAnimationFrame(animate);\n'
        '       if (isPlaying) { /* update animations */ }\n'
        '       controls.update();\n'
        '       renderer.render(scene, camera);\n'
        '   }\n'
        '3. 初始化完成后先渲染一帧再隐藏 loading:\n'
        '   renderer.render(scene, camera);\n'
        '   document.getElementById(“loading”).style.display = “none”;\n'
        '   animate();\n'
        '4. Raycaster 需要递归收集所有 Mesh (包括 Group 子对象):\n'
        '   const allMeshes = [];\n'
        '   scene.traverse(child => { if (child.isMesh) allMeshes.push(child); });\n'
        '5. objects 字典中只存可独立控制的 Mesh/Line/Points，不要存 Group\n\n'
        '## 容错要求 (必须)\n'
        '- 检测 WebGL 支持，不支持时显示友好提示\n'
        '- 初始化用 try-catch 包裹，失败时显示错误信息和重试按钮\n'
        '- 加载时显示 loading 遮罩\n\n'
        '只输出 HTML，不要任何解释。'
    )

    raw = _qwen_chat_raw_text(
        system_prompt,
        user_prompt,
        max_tokens=9000,
        temperature=0.3,
        timeout_seconds=get_settings().qwen_classroom_timeout_seconds,
        user=user,
    )

    # Strip markdown code blocks and extract HTML
    html = _extract_html_from_response(raw)

    # Derive title from item
    title = f"{item.title} — 3D 互动演示"

    return html, title


def _extract_html_from_response(text: str) -> str:
    """Extract HTML document from LLM response, stripping markdown wrappers."""
    # Remove markdown code block markers
    text = text.strip()
    if text.startswith("```html"):
        text = text[7:]
    elif text.startswith("```"):
        text = text[3:]
    if text.endswith("```"):
        text = text[:-3]
    text = text.strip()

    # If it already starts with DOCTYPE or html tag, return as-is
    if text.startswith("<!DOCTYPE") or text.startswith("<html"):
        return text

    # Try to find HTML document boundaries
    doctype_idx = text.find("<!DOCTYPE html>")
    html_start = text.find("<html")
    start = doctype_idx if doctype_idx != -1 else html_start
    if start != -1:
        html_end = text.rfind("</html>")
        if html_end != -1:
            return text[start:html_end + 7]

    # Fallback: wrap in basic HTML structure
    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head><meta charset="utf-8"/><meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>3D 互动演示</title></head>
<body>{text}</body>
</html>"""


def _write_raw_visualization_html(session_id: int, item: LearningSyllabusItem, html_content: str) -> Path:
    """Write a pre-generated HTML string directly to a file."""
    output_dir = Path(__file__).resolve().parents[2] / "generated" / "visualizations"
    output_dir.mkdir(parents=True, exist_ok=True)
    path = output_dir / f"classroom-{session_id}-item-{item.id}-visualization3d.html"
    path.write_text(html_content, encoding="utf-8")
    return path


def generate_classroom_visualization(
    db: Session,
    user: User,
    session_id: int,
    request: ClassroomVisualizationGenerateRequest,
) -> ClassroomSession:
    session = _get_session(db, user, session_id)
    item = _item_or_404(db, user, session.syllabus_item_id)
    project = _project_or_404(db, user, session.project_id)
    package = _classroom_package_or_error(session)

    is_3d = _normalize_widget_type(request.preferred_kind, allow_auto=True) == "visualization3d"

    if is_3d:
        html_content, demo_title = _generate_visualization3d_html_direct(
            project, item, package, request.instruction, user=user
        )
        html_path = _write_raw_visualization_html(session.id, item, html_content)
        content_data = {
            "widget_type": "visualization3d",
            "title": demo_title,
            "html": html_content,
        }
        source = "THU-MAIC/OpenMAIC direct HTML; widget_type=visualization3d"
    else:
        demo = _generate_visualization_demo(project, item, package, request.instruction, request.preferred_kind, user=user)
        html_path = _write_visualization_html(session.id, item, demo)
        content_data = demo.model_dump(mode="json")
        source = f"THU-MAIC/OpenMAIC scene widget; widget_type={demo.widget_type}"
        demo_title = demo.title

    resource = ClassroomResource(
        session_id=session.id,
        syllabus_item_id=item.id,
        project_id=project.id,
        user_id=user.id,
        resource_type="interactive_visualization",
        title=demo_title,
        content_data=content_data,
        file_path=str(html_path),
        source=source,
        status="ready",
    )
    db.add(resource)
    db.add(
        AgentTaskRecord(
            session_id=f"classroom-{session.id}",
            user_id=user.id,
            agent="SceneOutlineAgent+InteractiveWidgetRouter",
            status="completed",
            input_summary=f"{project.title} / {item.title}",
            output_summary=f"生成互动演示：{demo_title}",
            latency_ms=0,
        )
    )
    _write_event(
        db,
        project,
        user,
        "classroom_visualization_generated",
        f"生成互动演示：{demo_title}",
        {"session_id": session.id, "widget_type": content_data.get("widget_type", "")},
    )
    db.commit()
    return _get_session(db, user, session.id)


def _repair_visualization_demo_json(
    *,
    system_prompt: str,
    original_prompt: str,
    raw: Any,
    validation_error: str,
    widget_preference: str,
    user: Optional[User] = None,
) -> Any:
    """Give the LLM one chance to fix its broken visualization JSON, same pattern as classroom package repair."""
    repair_prompt = (
        "你刚才输出的互动演示 JSON 未通过系统结构校验。请只返回修复后的完整 JSON 对象，不要 Markdown，不要解释。\n"
        "修复原则：\n"
        "1. 必须保留原演示主题和教学意图，只补齐或改正不合格字段。\n"
        "2. 不得删除必填顶层字段：title, widget_type, demo_type, learning_goal, description, "
        "variables, nodes, edges, frames, controls, teaching_points, student_tasks, safety_notes, "
        "code_snippet, physics_scene。\n"
        "3. 如果 widget_type 是 visualization3d，则 physics_scene 必须是一个包含 scene_kind, gravity, camera, objects 的对象；"
        "每个 object 必须包含 id, label, role, shape, size, position, velocity, mass, color；"
        "frames 每帧必须包含 label, narrative, metrics（至少 force 和 speed 数值）。\n"
        "4. 如果 widget_type 不是 visualization3d，则 physics_scene 必须为 null；"
        "nodes 至少 4 个、edges 至少 3 个、frames 至少 4 帧。\n"
        '5. 所有数组不得为空；不要使用“待补充”“暂无”“N/A”“占位符”。\n'
        f"校验错误：{validation_error}\n"
        f"原始任务要求：{_truncate_for_prompt(original_prompt, 5000)}\n"
        f"用户 widget_type 偏好：{widget_preference}\n"
        f"待修复 JSON：{_truncate_for_prompt(json.dumps(raw, ensure_ascii=False), 9000)}\n"
    )
    return _qwen_chat_raw_json(
        system_prompt,
        repair_prompt,
        max_tokens=5200,
        temperature=0.1,
        timeout_seconds=get_settings().qwen_classroom_timeout_seconds,
        user=user,
    )


def _generate_visualization_demo(
    project: LearningProject,
    item: LearningSyllabusItem,
    package: _ClassroomPackage,
    instruction: str,
    preferred_kind: str = "auto",
    *,
    user: Optional[User] = None,
) -> _VisualizationDemo:
    widget_preference = _normalize_widget_type(preferred_kind, allow_auto=True)
    knowledge_context = build_rag_context_for_classroom(project, item, instruction)

    viz_system_prompt = (
        "你是 OpenMAIC 风格的 SceneOutlineAgent + InteractiveWidgetRouter。"
        "为当前课堂页生成一个可交互教学场景规格。"
        "先判断内容最适合 diagram、simulation、code 还是 visualization3d。"
        "只有真实空间结构、物理装置、几何、机械、分子、3D 坐标关系等内容才允许使用 visualization3d。"
        "只输出严格 JSON，不输出 Markdown、HTML 或 JavaScript。"
        f"{FORMULA_OUTPUT_INSTRUCTIONS}"
    )
    viz_user_prompt = (
        "输出顶层字段必须完整包含：title, widget_type, demo_type, learning_goal, description, "
        "variables, nodes, edges, frames, controls, teaching_points, student_tasks, safety_notes, "
        "code_snippet, physics_scene。\n"
        "widget_type 只能是 diagram, simulation, code, visualization3d。\n"
        f"用户偏好：{widget_preference}。如果为 auto，你必须根据主题自行选择最清楚的形态，不要默认 3D。\n"
        "diagram 适合概念关系、系统结构、论文脉络、课程知识框架；"
        "simulation 适合算法迭代、状态变化、数据流、协议交互、训练过程；"
        "code 适合 Python/工程复现/API/实验脚本讲解；"
        "visualization3d 只适合空间、物理、机械、几何、分子、传感器布置等确实需要三维理解的内容。\n"
        "demo_type 必须使用规范值，不得自造别名："
        "diagram=concept_map/system_diagram/flowchart/comparison_map；"
        "simulation=algorithm_trace/data_flow/process_simulation/state_machine；"
        "code=code_walkthrough/debug_trace/api_flow/reproduction_demo。"
        "算法执行、递归栈、迭代过程、复杂度对比必须用 widget_type=simulation 且 demo_type=algorithm_trace。\n"
        "当 widget_type 不是 visualization3d 时：physics_scene 必须为 null；nodes 至少 4 个；edges 至少 3 个；"
        "每个 node 必须包含 id, label, kind, x, y, color, detail；x/y 是 0-100 的数字；"
        "每个 edge 必须包含 source, target, label；source/target 必须引用已有 node id；"
        "frames 需要 4-8 帧，每帧必须包含 label, narrative, metrics, active_nodes, active_edges；"
        "active_nodes/active_edges 必须引用节点 id 和边 id/source-target。\n"
        "当 widget_type 是 code 时，code_snippet 必须给出与主题相关、可阅读的短代码；"
        "当 widget_type 是 visualization3d 时：physics_scene 必须包含 scene_kind, gravity, camera, objects, annotations；"
        "objects 4-9 个，每个必须包含 id, label, role, shape, size, position, velocity, mass, color, particle_emitter。"
        "role 可以是 source, emitter, target, relay, observer, sink, processor, storage, controller, sensor；"
        "particle_emitter 为 true 时对象会发射粒子拖尾，适合 role=source/emitter/sensor 的对象。"
        "scene_kind 必须是 signal_wave, network_packet, neural_activation, optimization_landscape, sorting_collision, graph_diffusion, general_physics, physics_system 之一。"
        "camera 必须包含 position 和 target 两个 [x,y,z] 数组。"
        "annotations 数组可为空或包含标注箭头对象 {origin:[x,y,z], direction:[x,y,z], length:float, color:'#hex'}。"
        "每个 frame 必须包含 label, narrative, metrics（至少含 force 和 speed 数值）；"
        "可选包含 camera_position 和 camera_target，用于切换帧时自动运镜过渡。\n"
        "controls 至少 2 个，每个包含 name, label, min_value, max_value, default_value, description。\n"
        "teaching_points 至少 3 条，student_tasks 至少 2 条，safety_notes 至少 1 条。\n"
        "所有内容必须贴合课堂当前页和知识库来源，不得输出空数组、占位符、模板化泛泛描述。\n"
        f"项目：{project.title}\n"
        f"研究/学习方向：{project.research_direction}\n"
        f"学习项：{item.title}\n"
        f"学习目标：{item.objective}\n"
        f"知识点：{item.knowledge_points}\n"
        f"课堂摘要：{package.learning_summary}\n"
        f"课程知识库来源：\n{knowledge_context}\n"
        f"当前页/用户要求：{instruction}\n"
    )

    raw = _qwen_chat_raw_json(viz_system_prompt, viz_user_prompt, user=user)
    try:
        normalized = _normalize_visualization_demo(raw, widget_preference)
        return _VisualizationDemo.model_validate(normalized)
    except LLMResponseError as first_error:
        repaired = _repair_visualization_demo_json(
            system_prompt=viz_system_prompt,
            original_prompt=viz_user_prompt,
            raw=raw,
            validation_error=str(first_error),
            widget_preference=widget_preference,
            user=user,
        )
        try:
            normalized = _normalize_visualization_demo(repaired, widget_preference)
            return _VisualizationDemo.model_validate(normalized)
        except (LLMResponseError, Exception) as second_error:
            raise LLMResponseError(
                "互动演示 JSON 自动修复后仍未通过校验："
                f"{second_error}；首次错误：{first_error}。请您截图联系管理员处理"
            ) from second_error


def _synthesize_physics_scene(
    raw_nodes: list[dict[str, Any]],
    raw_frames: list[dict[str, Any]],
    title: str,
) -> dict[str, Any]:
    """Build a minimal valid physics_scene from diagram-style nodes/frames data."""
    objects: list[dict[str, Any]] = []
    colors = _VISUAL_COLOR_PALETTE
    roles = ["source", "relay", "processor", "target", "observer", "sink"]
    for i, node in enumerate(raw_nodes[:8]):
        if not isinstance(node, dict):
            continue
        label = str(node.get("label") or node.get("title") or f"对象{i + 1}")
        node_id = str(node.get("id") or f"obj_{i + 1}")
        x = float(node.get("x", 50)) / 100 * 10 - 5  # 0-100 → -5..5
        z = float(node.get("y", 50)) / 100 * 10 - 5
        y = 1.5 + i * 0.6
        objects.append({
            "id": re.sub(r"[^a-zA-Z0-9_-]", "_", node_id)[:48],
            "label": str(label)[:30],
            "role": node.get("kind") or roles[i % len(roles)],
            "shape": "sphere",
            "size": [0.8, 0.8, 0.8],
            "position": [round(x, 2), round(y, 2), round(z, 2)],
            "velocity": [round((i - 3) * 0.2, 2), 0.0, round((i % 3 - 1) * 0.3, 2)],
            "mass": 1.0,
            "color": str(node.get("color") or colors[i % len(colors)]),
            "particle_emitter": i < 2,
        })
    # Ensure at least 4 objects
    while len(objects) < 4:
        i = len(objects)
        objects.append({
            "id": f"synth_obj_{i + 1}",
            "label": f"节点 {i + 1}",
            "role": roles[i % len(roles)],
            "shape": "sphere",
            "size": [0.7, 0.7, 0.7],
            "position": [(i - 2) * 2.0, 2.0, (i % 2) * 2.0 - 1.0],
            "velocity": [0.0, 0.0, 0.0],
            "mass": 1.0,
            "color": colors[i % len(colors)],
            "particle_emitter": False,
        })

    return {
        "scene_kind": "general_physics",
        "gravity": [0, -9.8, 0],
        "camera": {"position": [8, 6, 12], "target": [0, 2, 0]},
        "objects": objects[:9],
        "annotations": [],
    }


def _fill_3d_frames(frames: list[dict[str, Any]], title: str) -> list[dict[str, Any]]:
    """Fill missing 3D frames with defaults so we always have at least 4."""
    result = list(frames)
    while len(result) < 4:
        i = len(result) + 1
        result.append({
            "label": f"步骤 {i}",
            "metrics": {"step": float(i), "force": 1.0, "speed": 1.0},
            "narrative": f"观察「{title}」在第 {i} 步的变化。",
        })
    return result[:12]


def _normalize_visualization_demo(raw: Any, preferred_kind: str = "auto") -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError("互动演示 JSON 顶层必须是对象")
    widget_type = _normalize_widget_type(raw.get("widget_type") or raw.get("type"), allow_auto=False)
    preferred = _normalize_widget_type(preferred_kind, allow_auto=True)
    if preferred != "auto" and widget_type != preferred:
        # Force user's preferred widget_type; also remap demo_type to a valid fallback
        raw["widget_type"] = preferred
        widget_type = preferred
        raw_demo = str(raw.get("demo_type") or "")
        if preferred == "visualization3d":
            # Map non-3D demo types to a sensible 3D fallback
            raw["demo_type"] = _DEMO_TO_3D_FALLBACK.get(raw_demo.strip().lower(), "physics_system")
        else:
            raw["demo_type"] = _3D_TO_NON_3D_FALLBACK.get(preferred, {}).get(
                raw_demo.strip().lower(),
                {"diagram": "concept_map", "simulation": "algorithm_trace", "code": "code_walkthrough"}.get(preferred, "concept_map"),
            )

    title = _first_visual_text(raw.get("title") or raw.get("name"), "互动演示")
    learning_goal = _first_visual_text(raw.get("learning_goal") or raw.get("goal"), f"理解「{title}」的关键机制")
    description = _first_visual_text(raw.get("description") or raw.get("summary"), learning_goal)

    variables: list[str] = []
    for value in _as_list(raw.get("variables")):
        try:
            variables.append(_normalize_visualization_variable(value))
        except LLMResponseError:
            continue
    variables = [value for value in variables if value]
    if len(variables) < 2:
        for fallback in ("步骤进度", "当前状态", "关键指标", "理解程度"):
            if fallback not in variables:
                variables.append(fallback)
            if len(variables) >= 2:
                break

    controls: list[dict[str, Any]] = []
    for index, control in enumerate(_as_list(raw.get("controls")), start=1):
        try:
            controls.append(_normalize_visualization_control(control, index))
        except LLMResponseError:
            continue
    if len(controls) < 2:
        controls = _fill_visualization_controls(controls)

    safety_notes = _as_str_list(raw.get("safety_notes"))
    if not safety_notes:
        safety_notes = ["演示为概念可视化，结论需回到课程材料和来源文献核对。"]
    teaching_points = _as_str_list(raw.get("teaching_points"))
    if len(teaching_points) < 3:
        teaching_points = _fill_visualization_texts(
            teaching_points,
            [learning_goal, description, title],
            3,
            "观察当前步骤中被高亮的节点和关系。",
        )
    student_tasks = _as_str_list(raw.get("student_tasks"))
    if len(student_tasks) < 2:
        student_tasks = _fill_visualization_texts(
            student_tasks,
            [f"按顺序播放演示，复述「{title}」的关键变化。", "暂停在关键步骤，解释高亮节点之间的关系。"],
            2,
            "用一句话总结本演示的核心机制。",
        )

    if widget_type == "visualization3d":
        # If LLM returned non-3D data but we forced widget_type, synthesize a 3D scene
        raw_physics = raw.get("physics_scene")
        if not isinstance(raw_physics, dict) or not raw_physics.get("objects"):
            raw_nodes = _as_list(raw.get("nodes"))
            raw_frames = _as_list(raw.get("frames"))
            raw_physics = _synthesize_physics_scene(raw_nodes, raw_frames, title)
            raw["physics_scene"] = raw_physics

        demo_type = _normalize_demo_type(_require_text(raw.get("demo_type"), "visualization.demo_type"))
        frames: list[dict[str, Any]] = []
        for index, frame in enumerate(_as_list(raw.get("frames")), start=1):
            try:
                frames.append(_normalize_visualization_frame(frame, index, variables))
            except LLMResponseError:
                continue
        if len(frames) < 4:
            frames = _fill_3d_frames(frames, title)
        physics_scene = _normalize_physics_scene(raw_physics, demo_type)
        nodes: list[dict[str, Any]] = []
        edges: list[dict[str, str]] = []
    else:
        demo_type = _normalize_non_3d_demo_type(widget_type, str(raw.get("demo_type") or ""))
        nodes: list[dict[str, Any]] = []
        for index, node in enumerate(_as_list(raw.get("nodes")), start=1):
            try:
                nodes.append(_normalize_visual_node(node, index))
            except LLMResponseError:
                continue
        if len(nodes) < 4:
            nodes = _fill_visual_nodes(nodes, title, description, teaching_points, widget_type)
        node_ids = {node["id"] for node in nodes}
        edges: list[dict[str, str]] = []
        for index, edge in enumerate(_as_list(raw.get("edges")), start=1):
            try:
                edges.append(_normalize_visual_edge(edge, index, node_ids))
            except LLMResponseError:
                continue
        if len(edges) < 3:
            edges = _fill_visual_edges(edges, nodes)
        edge_ids = {edge["id"] for edge in edges}
        frames: list[dict[str, Any]] = []
        for index, frame in enumerate(_as_list(raw.get("frames")), start=1):
            try:
                frames.append(_normalize_adaptive_visual_frame(frame, index, variables, node_ids, edge_ids))
            except LLMResponseError:
                continue
        if len(frames) < 4:
            frames = _fill_visual_frames(frames, nodes, edges, title)
        physics_scene = None

    code_snippet = str(raw.get("code_snippet") or raw.get("code") or "").strip()
    if widget_type == "code" and len(code_snippet) < 40:
        code_snippet = (
            f"# {title}\n"
            "def observe_step(step):\n"
            "    return {'step': step, 'progress': step / 4}\n\n"
            "for step in range(1, 5):\n"
            "    print(observe_step(step))\n"
        )

    return {
        "title": title,
        "demo_type": demo_type,
        "widget_type": widget_type,
        "learning_goal": learning_goal,
        "description": description,
        "variables": variables[:8],
        "frames": frames[:12],
        "controls": controls,
        "teaching_points": teaching_points[:8],
        "student_tasks": student_tasks[:6],
        "safety_notes": safety_notes,
        "nodes": nodes[:12],
        "edges": edges[:18],
        "code_snippet": code_snippet,
        "physics_scene": physics_scene,
    }


def _normalize_widget_type(value: Any, allow_auto: bool) -> str:
    text = str(value or "auto").strip().lower().replace("-", "_").replace(" ", "_")
    aliases = {
        "3d": "visualization3d",
        "visualization_3d": "visualization3d",
        "visualisation3d": "visualization3d",
        "graph": "diagram",
        "mindmap": "diagram",
        "flow": "simulation",
        "sim": "simulation",
        "coding": "code",
    }
    text = aliases.get(text, text)
    allowed = set(_VISUAL_WIDGET_TYPES)
    if allow_auto:
        allowed.add("auto")
    if text not in allowed:
        raise LLMResponseError(f"互动演示 widget_type 非法：{value}；允许值：{', '.join(sorted(allowed))}")
    return text


def _normalize_non_3d_demo_type(widget_type: str, value: str) -> str:
    normalized = value.strip().lower().replace("-", "_").replace(" ", "_")
    normalized = _DEMO_TYPE_ALIASES.get(widget_type, {}).get(normalized, normalized)
    allowed = _NON_3D_DEMO_TYPES[widget_type]
    if normalized not in allowed:
        fallback = {
            "diagram": "concept_map",
            "simulation": "algorithm_trace",
            "code": "code_walkthrough",
        }.get(widget_type)
        if fallback:
            return fallback
        raise LLMResponseError(f"{widget_type} 演示 demo_type 非法：{value}；允许值：{', '.join(sorted(allowed))}")
    return normalized


def _fill_visualization_controls(controls: list[dict[str, Any]]) -> list[dict[str, Any]]:
    result = list(controls)
    defaults = [
        {
            "name": "playback_speed",
            "label": "播放速度",
            "min_value": 0.5,
            "max_value": 2.0,
            "default_value": 1.0,
            "description": "控制演示自动播放速度。",
        },
        {
            "name": "focus_strength",
            "label": "高亮强度",
            "min_value": 0.0,
            "max_value": 1.0,
            "default_value": 0.7,
            "description": "用于提示当前步骤的重点关系。",
        },
    ]
    existing = {str(control.get("name") or "") for control in result}
    for default in defaults:
        if len(result) >= 2:
            break
        if default["name"] not in existing:
            result.append(default)
    return result[:4]


def _fill_visualization_texts(values: list[str], candidates: list[str], min_count: int, default: str) -> list[str]:
    result = [_compact_text(value, 160) for value in values if _compact_text(value, 160)]
    for candidate in candidates:
        text = _compact_text(str(candidate or ""), 160)
        if text and text not in result:
            result.append(text)
        if len(result) >= min_count:
            return result[:min_count]
    while len(result) < min_count:
        result.append(default)
    return result[:min_count]


def _first_visual_text(raw: Any, fallback: str) -> str:
    texts = _as_str_list(raw)
    if texts:
        text = _compact_text(texts[0], 120)
        if text:
            return text
    return _compact_text(fallback, 120)


def _fill_visual_nodes(
    nodes: list[dict[str, Any]],
    title: str,
    description: str,
    teaching_points: list[str],
    widget_type: str,
) -> list[dict[str, Any]]:
    result = list(nodes)
    seeds = [title, *teaching_points, description, "输入", "处理", "观察", "输出"]
    positions = [(14, 50), (36, 24), (62, 32), (82, 58), (58, 78), (28, 76)]
    existing = {node["id"] for node in result}
    kind = "step" if widget_type in {"simulation"} else "concept"
    for seed in seeds:
        if len(result) >= 4:
            break
        label = _compact_text(str(seed or ""), 40) or f"节点 {len(result) + 1}"
        node_id = re.sub(r"[^a-zA-Z0-9_-]", "_", f"auto_{len(result) + 1}")[:48]
        while node_id in existing:
            node_id = f"{node_id}_{len(result) + 1}"
        x, y = positions[len(result) % len(positions)]
        result.append(
            {
                "id": node_id,
                "label": label,
                "kind": kind,
                "x": float(x),
                "y": float(y),
                "color": _VISUAL_COLOR_PALETTE[len(result) % len(_VISUAL_COLOR_PALETTE)],
                "detail": _compact_text(description or label, 140),
            }
        )
        existing.add(node_id)
    return result[:12]


def _fill_visual_edges(edges: list[dict[str, str]], nodes: list[dict[str, Any]]) -> list[dict[str, str]]:
    result = list(edges)
    existing = {edge["id"] for edge in result}
    for index in range(max(0, len(nodes) - 1)):
        if len(result) >= 3:
            break
        source = nodes[index]["id"]
        target = nodes[index + 1]["id"]
        edge_id = f"{source}->{target}"
        if edge_id in existing:
            continue
        result.append({"id": edge_id, "source": source, "target": target, "label": "推进"})
        existing.add(edge_id)
    return result[:18]


def _fill_visual_frames(
    frames: list[dict[str, Any]],
    nodes: list[dict[str, Any]],
    edges: list[dict[str, str]],
    title: str,
) -> list[dict[str, Any]]:
    result = list(frames)
    if not nodes:
        return result
    while len(result) < 4:
        index = len(result)
        node = nodes[min(index, len(nodes) - 1)]
        edge = edges[min(index, len(edges) - 1)] if edges else None
        result.append(
            {
                "label": f"步骤 {index + 1}",
                "metrics": {"progress": round((index + 1) / 4, 2), "step": float(index + 1)},
                "narrative": f"围绕「{title}」观察 {node['label']} 的作用。",
                "active_nodes": [node["id"]],
                "active_edges": [edge["id"]] if edge else [],
                "annotations": [],
            }
        )
    return result[:12]


def _normalize_visual_node(raw: Any, index: int) -> dict[str, Any]:
    if not isinstance(raw, dict):
        raise LLMResponseError(f"互动演示 nodes[{index}] 必须是对象")
    node_id = _require_text(raw.get("id"), f"visualization.nodes[{index}].id")
    return {
        "id": re.sub(r"[^a-zA-Z0-9_-]", "_", node_id)[:48],
        "label": _require_text(raw.get("label") or raw.get("title"), f"visualization.nodes[{index}].label")[:40],
        "kind": _require_text(raw.get("kind") or raw.get("type") or "concept", f"visualization.nodes[{index}].kind")[:32],
        "x": _bounded_coordinate(raw.get("x"), f"visualization.nodes[{index}].x"),
        "y": _bounded_coordinate(raw.get("y"), f"visualization.nodes[{index}].y"),
        "color": _visual_color(raw.get("color"), index),
        "detail": _require_text(raw.get("detail") or raw.get("description"), f"visualization.nodes[{index}].detail")[:140],
    }


def _normalize_visual_edge(raw: Any, index: int, node_ids: set[str]) -> dict[str, str]:
    if not isinstance(raw, dict):
        raise LLMResponseError(f"互动演示 edges[{index}] 必须是对象")
    source = re.sub(r"[^a-zA-Z0-9_-]", "_", _require_text(raw.get("source"), f"visualization.edges[{index}].source"))[:48]
    target = re.sub(r"[^a-zA-Z0-9_-]", "_", _require_text(raw.get("target"), f"visualization.edges[{index}].target"))[:48]
    if source not in node_ids or target not in node_ids:
        raise LLMResponseError(f"互动演示 edges[{index}] 引用了不存在的 node：{source}->{target}")
    edge_id = str(raw.get("id") or f"{source}->{target}").strip()
    return {
        "id": re.sub(r"[^a-zA-Z0-9_>.-]", "_", edge_id)[:96],
        "source": source,
        "target": target,
        "label": _require_text(raw.get("label") or raw.get("relation"), f"visualization.edges[{index}].label")[:48],
    }


def _normalize_adaptive_visual_frame(
    raw: Any,
    index: int,
    variables: list[str],
    node_ids: set[str],
    edge_ids: set[str],
) -> dict:
    if not isinstance(raw, dict):
        raise LLMResponseError(f"互动演示 frames[{index}] 必须是对象")
    metrics = _numeric_metrics(raw.get("metrics") if isinstance(raw.get("metrics"), dict) else {})
    if not metrics:
        metrics = {"step": float(index), "progress": min(1.0, index / 4)}
    active_nodes = _normalize_frame_targets(raw.get("active_nodes") or raw.get("nodes"), node_ids, f"frames[{index}].active_nodes")
    active_edges = _normalize_frame_targets(raw.get("active_edges") or raw.get("edges"), edge_ids, f"frames[{index}].active_edges")
    if not active_nodes:
        active_nodes = [next(iter(node_ids))]
    return {
        "label": _require_text(raw.get("label"), f"visualization.frames[{index}].label"),
        "metrics": metrics,
        "narrative": _require_text(raw.get("narrative"), f"visualization.frames[{index}].narrative"),
        "active_nodes": active_nodes,
        "active_edges": active_edges,
        "annotations": _as_str_list(raw.get("annotations"))[:4],
    }


def _normalize_frame_targets(raw: Any, allowed: set[str], field_path: str) -> list[str]:
    targets: list[str] = []
    for value in _as_list(raw):
        text = str(value).strip()
        if not text:
            continue
        normalized = re.sub(r"[^a-zA-Z0-9_>.-]", "_", text)[:96]
        if normalized not in allowed:
            continue
        targets.append(normalized)
    return targets


def _bounded_coordinate(value: Any, field_path: str) -> float:
    metric = _metric_value(value)
    if metric is None:
        raise LLMResponseError(f"互动演示 {field_path} 必须是 0-100 的数字")
    if 0 <= metric <= 1:
        metric *= 100
    return max(0.0, min(100.0, float(metric)))


def _visual_color(value: Any, index: int = 1) -> str:
    if value is None:
        return _VISUAL_COLOR_PALETTE[(index - 1) % len(_VISUAL_COLOR_PALETTE)]
    color = str(value).strip()
    if not re.fullmatch(r"#[0-9a-fA-F]{6}", color):
        return _VISUAL_COLOR_PALETTE[(index - 1) % len(_VISUAL_COLOR_PALETTE)]
    return color


def _write_visualization_html(session_id: int, item: LearningSyllabusItem, demo: _VisualizationDemo) -> Path:
    output_dir = Path(__file__).resolve().parents[2] / "generated" / "visualizations"
    output_dir.mkdir(parents=True, exist_ok=True)
    path = output_dir / f"classroom-{session_id}-item-{item.id}-visualization.html"
    if demo.widget_type == "visualization3d":
        html_doc = render_three_physics_html(demo.model_dump(mode="json"), demo.title)
    else:
        html_doc = render_adaptive_visualization_html(demo.model_dump(mode="json"), demo.title)
    path.write_text(html_doc, encoding="utf-8")
    return path
