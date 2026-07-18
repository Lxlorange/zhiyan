from __future__ import annotations

import hashlib
import io
import json
import os
import re
import shutil
import subprocess
import tempfile
import threading
import zipfile
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Optional

from fastapi import UploadFile
from pydantic import AliasChoices, BaseModel, ConfigDict, Field
from pypdf import PdfReader
from sqlalchemy import delete, func, select, text
from sqlalchemy.orm import Session

from app.core.database import SessionLocal
from app.core.config import get_settings
from app.models.learning import (
    Course,
    CourseChapter,
    DocumentChunk,
    KnowledgeDocument,
    KnowledgeImportJob,
    KnowledgePoint,
)
from app.models.user import User
from app.schemas import KnowledgeChunkRead, KnowledgeDocumentRead
from app.services.llm_client import LLMResponseError, qwen_chat_json, validate_qwen_config


SUPPORTED_TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".py", ".java", ".ts", ".js", ".html", ".css"}
SUPPORTED_OFFICE_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".docx", ".doc"}
SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}
SUPPORTED_ARCHIVES = {".zip"}
SUPPORTED_EXTENSIONS = SUPPORTED_TEXT_EXTENSIONS | SUPPORTED_OFFICE_EXTENSIONS | SUPPORTED_IMAGE_EXTENSIONS
MAX_SINGLE_EXTRACTED_CHARS = 1_000_000
CONTROL_CHARS_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")


class KnowledgeIngestionError(RuntimeError):
    status_code: int

    def __init__(self, message: str, status_code: int = 422):
        super().__init__(message)
        self.status_code = status_code


class KnowledgeImportJobRead(BaseModel):
    id: int
    course_code: str
    course_title: str
    source_name: str
    status: str
    total_files: int
    parsed_files: int
    failed_files: int
    total_chunks: int
    error_message: str
    options: dict = Field(default_factory=dict)
    created_at: datetime
    updated_at: datetime

    model_config = {"from_attributes": True}


class KnowledgeStorageUsageRead(BaseModel):
    quota_bytes: int
    used_bytes: int
    remaining_bytes: int
    quota_mb: int
    used_mb: float
    remaining_mb: float
    used_percent: float
    document_count: int
    job_count: int


class ImportedKnowledgeDependency(BaseModel):
    model_config = ConfigDict(populate_by_name=True)

    topic_name: str = Field(validation_alias=AliasChoices("topic_name", "topicName", "topicId"))
    prerequisite_name: str = Field(validation_alias=AliasChoices("prerequisite_name", "prerequisiteName", "prerequisiteId"))
    strength: str = Field(default="soft", pattern="^(hard|soft)$")
    reason: str = ""


class ImportedKnowledgeTopic(BaseModel):
    model_config = ConfigDict(populate_by_name=True)

    id: str = Field(default="", validation_alias=AliasChoices("id", "topicId"))
    type: str = Field(default="CONCEPTUAL", pattern="^(CONCEPTUAL|PROCEDURAL|REPRESENTATIONAL|LANGUAGE|META)$")
    subject: str = ""
    domain: str = ""
    name: str
    description: str
    evidence: list[str] = Field(default_factory=list)
    assessment_prompt: str = Field(default="", validation_alias=AliasChoices("assessment_prompt", "assessmentPrompt"))
    standards: list[str] = Field(default_factory=list)
    source_cues: list[str] = Field(default_factory=list, validation_alias=AliasChoices("source_cues", "sourceCues"))


class ImportedKnowledgeTaxonomy(BaseModel):
    model_config = ConfigDict(populate_by_name=True)

    topics: list[ImportedKnowledgeTopic] = Field(default_factory=list)
    dependencies: list[ImportedKnowledgeDependency] = Field(default_factory=list)


@dataclass
class ParsedSection:
    text: str
    section_title: str = ""
    page_no: Optional[int] = None
    slide_no: Optional[int] = None
    metadata: Optional[dict] = None


@dataclass
class ParsedKnowledgeFile:
    path: Path
    title: str
    doc_type: str
    source_uri: str
    file_hash: str
    sections: list[ParsedSection]
    metadata: dict


def list_import_jobs(db: Session, user: User, limit: int = 20) -> list[KnowledgeImportJobRead]:
    jobs = db.scalars(
        select(KnowledgeImportJob)
        .where(KnowledgeImportJob.user_id == user.id)
        .order_by(KnowledgeImportJob.created_at.desc())
        .limit(limit)
    ).all()
    return [KnowledgeImportJobRead.model_validate(job) for job in jobs]


def get_import_job(db: Session, user: User, job_id: int) -> KnowledgeImportJobRead:
    job = db.get(KnowledgeImportJob, job_id)
    if job is None or job.user_id != user.id:
        raise KnowledgeIngestionError("知识库导入任务不存在", 404)
    return KnowledgeImportJobRead.model_validate(job)


def delete_import_job(db: Session, user: User, job_id: int) -> None:
    job = db.get(KnowledgeImportJob, job_id)
    if job is None or job.user_id != user.id:
        raise KnowledgeIngestionError("知识库导入任务不存在", 404)
    if job.status in {"queued", "running"}:
        raise KnowledgeIngestionError("资料仍在后台解析中，请等待完成或失败后再删除。", 409)
    for document in _documents_for_job(db, job_id):
        db.delete(document)
    _delete_job_storage_dir(job.id)
    db.delete(job)
    db.commit()


def get_knowledge_storage_usage(db: Session, user: User) -> KnowledgeStorageUsageRead:
    settings = get_settings()
    quota_bytes = settings.knowledge_storage_quota_mb * 1024 * 1024
    job_ids = db.scalars(select(KnowledgeImportJob.id).where(KnowledgeImportJob.user_id == user.id)).all()
    used_bytes = sum(_job_storage_bytes(job_id) for job_id in job_ids)
    document_count = sum(len(_documents_for_job(db, job_id)) for job_id in job_ids)
    remaining_bytes = max(0, quota_bytes - used_bytes)
    return KnowledgeStorageUsageRead(
        quota_bytes=quota_bytes,
        used_bytes=used_bytes,
        remaining_bytes=remaining_bytes,
        quota_mb=settings.knowledge_storage_quota_mb,
        used_mb=round(used_bytes / 1024 / 1024, 2),
        remaining_mb=round(remaining_bytes / 1024 / 1024, 2),
        used_percent=round(min(100, used_bytes / quota_bytes * 100), 1) if quota_bytes else 0,
        document_count=document_count,
        job_count=len(job_ids),
    )


def mark_interrupted_import_jobs(db: Session) -> int:
    jobs = db.scalars(
        select(KnowledgeImportJob).where(KnowledgeImportJob.status.in_(["queued", "running"]))
    ).all()
    for job in jobs:
        job.status = "failed"
        job.error_message = job.error_message or "后端服务在导入过程中停止，任务已中断，请重新上传。"
        _set_import_job_progress(job, status="failed", stage="导入已中断", percent=100, active_file="")
    if jobs:
        db.commit()
    return len(jobs)


def list_knowledge_documents(
    db: Session,
    user: User,
    *,
    course_code: Optional[str] = None,
    query: str = "",
    limit: int = 50,
) -> list[KnowledgeDocumentRead]:
    limit = max(1, min(limit, 200))
    statement = (
        select(KnowledgeDocument, func.count(DocumentChunk.id).label("chunk_count"))
        .outerjoin(DocumentChunk, DocumentChunk.document_id == KnowledgeDocument.id)
        .group_by(KnowledgeDocument.id)
        .order_by(KnowledgeDocument.created_at.desc(), KnowledgeDocument.id.desc())
        .limit(limit)
    )
    if course_code:
        statement = statement.where(KnowledgeDocument.course_code == course_code)
    if query:
        like = f"%{query.strip()}%"
        statement = statement.where(
            (KnowledgeDocument.title.ilike(like))
            | (KnowledgeDocument.file_name.ilike(like))
            | (KnowledgeDocument.summary.ilike(like))
            | (KnowledgeDocument.course_code.ilike(like))
        )
    rows = db.execute(statement).all()
    return [
        _document_read(document, chunk_count)
        for document, chunk_count in rows
        if _user_can_manage_document(db, user, document)
    ]


def list_document_chunks(db: Session, user: User, document_id: int, *, limit: int = 100) -> list[KnowledgeChunkRead]:
    document = _get_manageable_document(db, user, document_id)
    limit = max(1, min(limit, 300))
    chunks = db.execute(
        select(DocumentChunk, KnowledgePoint.name)
        .join(KnowledgePoint, KnowledgePoint.id == DocumentChunk.knowledge_point_id, isouter=True)
        .where(DocumentChunk.document_id == document.id)
        .order_by(DocumentChunk.chunk_index)
        .limit(limit)
    ).all()
    return [
        KnowledgeChunkRead(
            id=chunk.id,
            document_id=chunk.document_id,
            chunk_index=chunk.chunk_index,
            content=chunk.content,
            keywords=list(chunk.keywords or []),
            knowledge_point=point_name or "课程资料",
            page_no=chunk.page_no,
            slide_no=chunk.slide_no,
            section_title=chunk.section_title,
            token_count=chunk.token_count,
        )
        for chunk, point_name in chunks
    ]


def delete_knowledge_document(db: Session, user: User, document_id: int) -> None:
    document = _get_manageable_document(db, user, document_id)
    db.delete(document)
    db.commit()


def _document_read(document: KnowledgeDocument, chunk_count: int) -> KnowledgeDocumentRead:
    return KnowledgeDocumentRead(
        id=document.id,
        title=document.title,
        doc_type=document.doc_type,
        source_uri=document.source_uri,
        summary=document.summary,
        file_name=document.file_name,
        file_hash=document.file_hash,
        course_code=document.course_code,
        parse_status=document.parse_status,
        parse_meta=document.parse_meta or {},
        chunk_count=int(chunk_count or 0),
        created_at=document.created_at,
    )


def _get_manageable_document(db: Session, user: User, document_id: int) -> KnowledgeDocument:
    document = db.get(KnowledgeDocument, document_id)
    if document is None:
        raise KnowledgeIngestionError("知识库文档不存在", 404)
    if not _user_can_manage_document(db, user, document):
        raise KnowledgeIngestionError("没有权限管理该知识库文档", 403)
    return document


def _user_can_manage_document(db: Session, user: User, document: KnowledgeDocument) -> bool:
    job_id = _job_id_from_document_path(document.file_path)
    if job_id is None:
        return False
    job = db.get(KnowledgeImportJob, job_id)
    return bool(job and job.user_id == user.id)


def _job_id_from_document_path(value: str) -> Optional[int]:
    match = re.search(r"(?:^|[\\/])job-(\d+)(?:[\\/]|$)", value or "")
    if not match:
        return None
    return int(match.group(1))


async def import_knowledge_upload(
    db: Session,
    user: User,
    file: UploadFile,
    *,
    course_code: str,
    course_title: str,
    use_ocr: bool,
    rebuild_course: bool,
) -> KnowledgeImportJobRead:
    if not file.filename:
        raise KnowledgeIngestionError("上传文件缺少文件名", 400)
    settings = get_settings()
    content = await file.read()
    if not content:
        raise KnowledgeIngestionError("上传文件为空", 400)
    max_bytes = settings.knowledge_max_upload_mb * 1024 * 1024
    if len(content) > max_bytes:
        raise KnowledgeIngestionError(f"上传文件过大，当前上限 {settings.knowledge_max_upload_mb} MB", 413)
    usage = get_knowledge_storage_usage(db, user)
    if usage.used_bytes + len(content) > usage.quota_bytes:
        remaining_mb = max(0, usage.remaining_bytes / 1024 / 1024)
        raise KnowledgeIngestionError(
            f"知识库空间不足。每位用户可用 {settings.knowledge_storage_quota_mb} MB，"
            f"当前已用 {usage.used_mb} MB，剩余 {remaining_mb:.2f} MB。请删除无用上传记录后再上传。",
            413,
        )

    upload_root = Path(settings.knowledge_upload_dir)
    upload_root.mkdir(parents=True, exist_ok=True)

    source_name = _safe_filename(file.filename)
    course_code = _clean_text(course_code)[:64] or "PERSONAL-KNOWLEDGE"
    course_title = _clean_text(course_title)[:255] or "个人知识库"
    job = KnowledgeImportJob(
        user_id=user.id,
        course_code=course_code,
        course_title=course_title,
        source_name=source_name,
        status="queued",
        options={
            "use_ocr": use_ocr,
            "rebuild_course": rebuild_course,
            "uploaded_bytes": len(content),
            "progress_stage": "文件已保存，等待后台解析",
            "progress_percent": 5,
        },
    )
    db.add(job)
    db.commit()
    db.refresh(job)

    job_dir = upload_root / f"job-{job.id}"
    job_dir.mkdir(parents=True, exist_ok=True)
    source_path = job_dir / source_name
    source_path.write_bytes(content)

    _start_import_job_worker(job.id)
    return KnowledgeImportJobRead.model_validate(job)


def _start_import_job_worker(job_id: int) -> None:
    worker = threading.Thread(
        target=_run_import_job_worker,
        args=(job_id,),
        name=f"knowledge-import-{job_id}",
        daemon=True,
    )
    worker.start()


def _run_import_job_worker(job_id: int) -> None:
    db = SessionLocal()
    try:
        _process_import_job(db, job_id)
    finally:
        db.close()


def _process_import_job(db: Session, job_id: int) -> None:
    job = db.get(KnowledgeImportJob, job_id)
    if job is None:
        return

    settings = get_settings()
    job_dir = Path(settings.knowledge_upload_dir) / f"job-{job.id}"
    source_path = job_dir / job.source_name
    use_ocr = bool((job.options or {}).get("use_ocr"))
    rebuild_course = bool((job.options or {}).get("rebuild_course"))
    course_code = job.course_code
    course_title = job.course_title

    try:
        _set_import_job_progress(job, status="running", stage="解析上传文件", percent=10)
        db.commit()
        db.refresh(job)

        parsed_files = _parse_upload_source(source_path, job_dir, use_ocr=use_ocr)
        if not parsed_files:
            raise KnowledgeIngestionError("上传资料中没有可解析的课件或文档", 422)

        total_chunks = 0
        job.total_files = len(parsed_files)
        _set_import_job_progress(job, stage="清洗资料并提取结构", percent=20)
        db.commit()
        db.refresh(job)
        course = _ensure_course(db, course_code, course_title)
        chapter = _ensure_import_chapter(db, course)
        point_cache = _load_point_cache(db, chapter.id)
        for file_index, parsed in enumerate(parsed_files, start=1):
            try:
                base_percent = 20 + int(65 * (file_index - 1) / max(1, len(parsed_files)))
                _set_import_job_progress(
                    job,
                    stage=f"抽取核心知识点：{_clean_text(parsed.title or parsed.path.name)[:80]}",
                    percent=max(25, base_percent),
                    active_file=parsed.path.name,
                )
                db.commit()
                db.refresh(job)
                if rebuild_course and total_chunks == 0 and job.parsed_files == 0 and job.failed_files == 0:
                    _clear_course_imports(db, course_code)
                    course = _ensure_course(db, course_code, course_title)
                    chapter = _ensure_import_chapter(db, course)
                    point_cache = _load_point_cache(db, chapter.id)
                chunks = _persist_parsed_file(db, course, chapter, point_cache, parsed)
                total_chunks += chunks
                job.parsed_files += 1
                job.total_chunks = total_chunks
                _set_import_job_progress(
                    job,
                    stage=f"已入库 {job.parsed_files}/{job.total_files} 个文件",
                    percent=20 + int(70 * file_index / max(1, len(parsed_files))),
                    active_file=parsed.path.name,
                )
                db.flush()
                db.commit()
                db.refresh(job)
            except Exception as exc:
                db.rollback()
                job = db.get(KnowledgeImportJob, job.id)
                if job is None:
                    raise KnowledgeIngestionError("知识库导入任务状态丢失，请重新上传", 500) from exc
                job.failed_files += 1
                job.error_message += f"{_clean_text(parsed.path.name)}: {_clean_text(str(exc))}\n"
                _set_import_job_progress(
                    job,
                    stage=f"文件入库失败：{_clean_text(parsed.path.name)[:80]}",
                    percent=20 + int(70 * file_index / max(1, len(parsed_files))),
                )
                db.commit()
                db.refresh(job)
        if job.parsed_files == 0:
            job.total_chunks = 0
            job.status = "failed"
            job.updated_at = datetime.utcnow()
            _set_import_job_progress(
                job,
                status="failed",
                stage="导入失败",
                percent=100,
            )
            db.commit()
            job.error_message = (
                "上传资料已解析，但全部文件入库失败：\n"
                + (job.error_message or "请检查文件内容和模型 Embedding 配置。")
            )
            db.commit()
            return
        job.total_chunks = total_chunks
        job.status = "completed" if job.failed_files == 0 else "partial_failed"
        job.updated_at = datetime.utcnow()
        _set_import_job_progress(
            job,
            status=job.status,
            stage="导入完成" if job.status == "completed" else "部分文件导入失败",
            percent=100,
            active_file="",
        )
        db.commit()
    except Exception as exc:
        db.rollback()
        job = db.get(KnowledgeImportJob, job.id)
        if job is not None:
            job.status = "failed"
            job.error_message = _clean_text(str(exc))
            job.updated_at = datetime.utcnow()
            _set_import_job_progress(job, status="failed", stage="导入失败", percent=100)
            db.commit()


def _set_import_job_progress(
    job: KnowledgeImportJob,
    *,
    status: Optional[str] = None,
    stage: Optional[str] = None,
    percent: Optional[int] = None,
    active_file: Optional[str] = None,
) -> None:
    options = dict(job.options or {})
    if status:
        job.status = status
    if stage is not None:
        options["progress_stage"] = stage
    if percent is not None:
        options["progress_percent"] = max(0, min(100, int(percent)))
    if active_file is not None:
        options["active_file"] = active_file
    options["progress_updated_at"] = datetime.utcnow().isoformat()
    job.options = options
    job.updated_at = datetime.utcnow()


def rebuild_missing_embeddings(db: Session, limit: int = 200) -> int:
    rows = db.scalars(
        select(DocumentChunk)
        .where(DocumentChunk.embedding.is_(None))
        .order_by(DocumentChunk.id)
        .limit(limit)
    ).all()
    for chunk in rows:
        chunk.embedding = embed_text(chunk.content)
    db.commit()
    return len(rows)


def _parse_upload_source(source_path: Path, job_dir: Path, *, use_ocr: bool) -> list[ParsedKnowledgeFile]:
    extension = source_path.suffix.lower()
    candidates: list[Path] = []
    if extension in SUPPORTED_ARCHIVES:
        extract_dir = job_dir / "extracted"
        extract_dir.mkdir(parents=True, exist_ok=True)
        _extract_zip(source_path, extract_dir)
        candidates = [path for path in extract_dir.rglob("*") if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS]
    elif extension in SUPPORTED_EXTENSIONS:
        candidates = [source_path]
    else:
        raise KnowledgeIngestionError(f"不支持的上传格式：{source_path.name}", 415)

    parsed: list[ParsedKnowledgeFile] = []
    errors: list[str] = []
    for candidate in candidates:
        try:
            parsed.append(_parse_single_file(candidate, use_ocr=use_ocr))
        except Exception as exc:
            errors.append(f"{candidate.name}: {exc}")
    if not parsed and errors:
        raise KnowledgeIngestionError("所有文件解析失败：\n" + "\n".join(errors[:20]), 422)
    return parsed


def _extract_zip(source_path: Path, extract_dir: Path) -> None:
    try:
        with zipfile.ZipFile(source_path) as archive:
            for member in archive.infolist():
                member_path = Path(_clean_zip_member_name(member.filename))
                if member.is_dir():
                    continue
                if member_path.is_absolute() or ".." in member_path.parts:
                    raise KnowledgeIngestionError(f"ZIP 包含不安全路径：{member.filename}", 400)
                target = extract_dir / member_path
                target.parent.mkdir(parents=True, exist_ok=True)
                with archive.open(member) as source, target.open("wb") as output:
                    shutil.copyfileobj(source, output)
    except zipfile.BadZipFile as exc:
        raise KnowledgeIngestionError(f"ZIP 文件损坏或格式不正确：{source_path.name}", 422) from exc


def _clean_zip_member_name(filename: str) -> str:
    parts = [_safe_filename(part) for part in Path(filename).parts if part not in {"", ".", ".."}]
    if not parts:
        return "upload.bin"
    return str(Path(*parts))


def _parse_single_file(path: Path, *, use_ocr: bool) -> ParsedKnowledgeFile:
    extension = path.suffix.lower()
    if extension == ".ppt":
        path = _convert_legacy_office(path, ".pptx")
        extension = path.suffix.lower()
    if extension == ".doc":
        path = _convert_legacy_office(path, ".docx")
        extension = path.suffix.lower()

    if extension == ".pdf":
        sections = _parse_pdf(path, use_ocr=use_ocr)
    elif extension == ".pptx":
        sections = _parse_pptx(path)
    elif extension == ".docx":
        sections = _parse_docx(path)
    elif extension in SUPPORTED_TEXT_EXTENSIONS:
        sections = _parse_text(path)
    elif extension in SUPPORTED_IMAGE_EXTENSIONS:
        sections = _parse_image(path, use_ocr=use_ocr)
    else:
        raise KnowledgeIngestionError(f"不支持的课件类型：{path.name}", 415)

    text_len = sum(len(section.text) for section in sections)
    if text_len <= 0:
        raise KnowledgeIngestionError(f"未提取到可用正文：{path.name}", 422)
    if text_len > MAX_SINGLE_EXTRACTED_CHARS:
        raise KnowledgeIngestionError(f"单个文件正文过长：{path.name}，请拆分后导入", 413)
    return ParsedKnowledgeFile(
        path=path,
        title=_clean_text(path.stem)[:255] or "导入资料",
        doc_type=extension.removeprefix("."),
        source_uri=_clean_text(str(path)),
        file_hash=_sha256_file(path),
        sections=sections,
        metadata={"extension": extension, "text_chars": text_len},
    )

def _parse_pdf(path: Path, *, use_ocr: bool) -> list[ParsedSection]:
    sections: list[ParsedSection] = []
    reader = PdfReader(str(path))
    for page_index, page in enumerate(reader.pages, start=1):
        text_value = _normalize_text(page.extract_text() or "")
        if text_value:
            sections.append(ParsedSection(text=text_value, page_no=page_index, section_title=f"page {page_index}"))
    if not sections and use_ocr:
        sections = _ocr_pdf_or_image(path)
    if not sections:
        raise KnowledgeIngestionError(f"PDF 未提取到文本：{path.name}。如为扫描版，请启用 OCR。", 422)
    return sections


def _parse_pptx(path: Path) -> list[ParsedSection]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    sections: list[ParsedSection] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text:
                parts.append(_clean_text(shape.text))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(_clean_text(cell.text) for cell in row.cells if cell.text))
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text
        except Exception:
            notes = ""
        if notes:
            parts.append("notes: " + _clean_text(notes))
        text_value = _normalize_text("\n".join(part.strip() for part in parts if part and part.strip()))
        if text_value:
            sections.append(ParsedSection(text=text_value, slide_no=slide_index, section_title=f"slide {slide_index}"))
    if not sections:
        raise KnowledgeIngestionError(f"PPTX 未提取到文本：{path.name}", 422)
    return sections


def _parse_docx(path: Path) -> list[ParsedSection]:
    from docx import Document

    document = Document(str(path))
    sections: list[ParsedSection] = []
    current_title = ""
    buffer: list[str] = []
    for paragraph in document.paragraphs:
        value = _normalize_text(paragraph.text)
        if not value:
            continue
        style_name = paragraph.style.name if paragraph.style else ""
        if style_name.lower().startswith("heading") and buffer:
            sections.append(ParsedSection(text="\n".join(buffer), section_title=current_title))
            buffer = []
            current_title = value
        else:
            buffer.append(value)
    table_text = []
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(_clean_text(cell.text).strip() for cell in row.cells if _clean_text(cell.text).strip())
            if row_text:
                table_text.append(row_text)
    buffer.extend(table_text)
    if buffer:
        sections.append(ParsedSection(text="\n".join(buffer), section_title=current_title))
    if not sections:
        raise KnowledgeIngestionError(f"DOCX 未提取到文本：{path.name}", 422)
    return sections


def _parse_text(path: Path) -> list[ParsedSection]:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            text_value = _normalize_text(raw.decode(encoding))
            break
        except UnicodeDecodeError:
            continue
    else:
        raise KnowledgeIngestionError(f"文本文件无法按 UTF-8 或 GB18030 解码：{path.name}", 422)
    if not text_value:
        raise KnowledgeIngestionError(f"文本文件为空：{path.name}", 422)
    return [ParsedSection(text=text_value, section_title=path.name)]


def _parse_image(path: Path, *, use_ocr: bool) -> list[ParsedSection]:
    if not use_ocr:
        raise KnowledgeIngestionError(f"图片文件需要启用 OCR 后才能导入：{path.name}", 422)
    return _ocr_pdf_or_image(path)


def _ocr_pdf_or_image(path: Path) -> list[ParsedSection]:
    settings = get_settings()
    if settings.local_ocr_engine.lower() != "tesseract":
        raise KnowledgeIngestionError(f"不支持的本地 OCR 引擎：{settings.local_ocr_engine}", 501)
    pytesseract, Image = _load_ocr_dependencies()
    if settings.tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd
    if path.suffix.lower() == ".pdf":
        return _ocr_pdf(path, pytesseract)
    text_value = _normalize_text(pytesseract.image_to_string(Image.open(path), lang="chi_sim+eng"))
    if not text_value:
        raise KnowledgeIngestionError(f"OCR 未识别到文字：{path.name}", 422)
    return [ParsedSection(text=text_value, section_title=path.name, metadata={"ocr": True})]


def _load_ocr_dependencies():
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise KnowledgeIngestionError("缺少 OCR 依赖，请安装 pytesseract 和 Pillow", 500) from exc
    return pytesseract, Image


def _ocr_pdf(path: Path, pytesseract) -> list[ParsedSection]:
    try:
        import fitz
    except ImportError as exc:
        raise KnowledgeIngestionError("扫描版 PDF OCR 缺少 PyMuPDF，请安装 backend/requirements.txt 后重试。", 500) from exc

    sections: list[ParsedSection] = []
    try:
        with fitz.open(str(path)) as document:
            for page_index, page in enumerate(document, start=1):
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                image_bytes = pixmap.tobytes("png")
                text_value = _normalize_text(
                    pytesseract.image_to_string(_image_from_bytes(image_bytes), lang="chi_sim+eng")
                )
                if text_value:
                    sections.append(
                        ParsedSection(
                            text=text_value,
                            page_no=page_index,
                            section_title=f"page {page_index} OCR",
                            metadata={"ocr": True},
                        )
                    )
    except KnowledgeIngestionError:
        raise
    except Exception as exc:
        raise KnowledgeIngestionError(f"PDF OCR 失败：{path.name}，{_clean_text(str(exc))}", 422) from exc
    if not sections:
        raise KnowledgeIngestionError(f"OCR 未识别到文字：{path.name}", 422)
    return sections


def _image_from_bytes(image_bytes: bytes):
    from PIL import Image

    return Image.open(io.BytesIO(image_bytes))


def _convert_legacy_office(path: Path, target_suffix: str) -> Path:
    settings = get_settings()
    libreoffice = settings.libreoffice_path.strip() or shutil.which("soffice") or shutil.which("libreoffice")
    if not libreoffice:
        raise KnowledgeIngestionError(f"{path.suffix} 旧格式需要安装 LibreOffice 后转换：{path.name}", 501)
    libreoffice_path = Path(libreoffice)
    if not libreoffice_path.exists():
        raise KnowledgeIngestionError(f"LIBREOFFICE_PATH 指向的文件不存在：{libreoffice}", 501)
    output_dir = path.parent / "converted"
    output_dir.mkdir(exist_ok=True)
    cmd = [str(libreoffice_path), "--headless", "--convert-to", target_suffix.removeprefix("."), "--outdir", str(output_dir), str(path)]
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=120)
    if result.returncode != 0:
        raise KnowledgeIngestionError(f"LibreOffice 转换失败：{path.name}，{result.stderr or result.stdout}", 422)
    converted = output_dir / f"{path.stem}{target_suffix}"
    if not converted.exists():
        raise KnowledgeIngestionError(f"LibreOffice 转换后未找到目标文件：{converted.name}", 422)
    return converted


def _ensure_course(db: Session, code: str, title: str) -> Course:
    course = db.scalar(select(Course).where(Course.code == code))
    if course is None:
        course = Course(code=code, title=title, description=f"{title} 课程知识库")
        db.add(course)
        db.flush()
    else:
        course.title = title
    return course


def _ensure_import_chapter(db: Session, course: Course) -> CourseChapter:
    chapter = db.scalar(
        select(CourseChapter).where(CourseChapter.course_id == course.id, CourseChapter.order_index == 999)
    )
    if chapter is None:
        chapter = CourseChapter(
            course_id=course.id,
            order_index=999,
            title=course.title,
            summary="由用户上传资料自动提取、去重并归并的核心知识点。",
        )
        db.add(chapter)
        db.flush()
    else:
        chapter.title = course.title
        chapter.summary = "由用户上传资料自动提取、去重并归并的核心知识点。"
    return chapter


def _load_point_cache(db: Session, chapter_id: int) -> dict[str, KnowledgePoint]:
    points = db.scalars(select(KnowledgePoint).where(KnowledgePoint.chapter_id == chapter_id)).all()
    return {_point_cache_key(point.name): point for point in points}


def _clear_course_imports(db: Session, course_code: str) -> None:
    docs = db.scalars(select(KnowledgeDocument).where(KnowledgeDocument.course_code == course_code)).all()
    for document in docs:
        db.delete(document)
    db.flush()


def _persist_parsed_file(
    db: Session,
    course: Course,
    chapter: CourseChapter,
    point_cache: dict[str, KnowledgePoint],
    parsed: ParsedKnowledgeFile,
) -> int:
    existing = db.scalar(select(KnowledgeDocument).where(KnowledgeDocument.file_hash == parsed.file_hash))
    title = _clean_text(parsed.title)[:255] or "导入资料"
    source_uri = _clean_text(parsed.source_uri)
    file_path = _clean_text(str(parsed.path))
    file_name = _safe_filename(parsed.path.name)
    summary = _summarize_sections(parsed.sections)
    parse_meta = _clean_metadata(parsed.metadata)
    chunks = list(_chunk_sections(parsed.sections))
    taxonomy = _build_import_taxonomy(title, parsed, chunks)
    parse_meta = {**parse_meta, **_taxonomy_parse_meta(taxonomy)}
    if existing is not None:
        db.execute(delete(DocumentChunk).where(DocumentChunk.document_id == existing.id))
        document = existing
        document.course_id = course.id
        document.title = title
        document.doc_type = parsed.doc_type
        document.source_uri = source_uri
        document.summary = summary
        document.file_path = file_path
        document.file_name = file_name
        document.course_code = course.code
        document.parse_status = "ready"
        document.parse_meta = parse_meta
    else:
        document = KnowledgeDocument(
            course_id=course.id,
            title=title,
            doc_type=parsed.doc_type,
            source_uri=source_uri,
            summary=summary,
            file_path=file_path,
            file_name=file_name,
            file_hash=parsed.file_hash,
            course_code=course.code,
            parse_status="ready",
            parse_meta=parse_meta,
        )
        db.add(document)
        db.flush()

    topic_points = _ensure_taxonomy_points(db, chapter, point_cache, taxonomy, title, chunks)
    for index, section in enumerate(chunks, start=1):
        safe_content = _clean_text(section.text)
        point = _match_chunk_to_topic(section, topic_points)
        topic = next((candidate for candidate, candidate_point in topic_points if candidate_point.id == point.id), None)
        chunk_meta = _clean_metadata(section.metadata or {})
        if isinstance(chunk_meta, dict) and topic is not None:
            chunk_meta = {
                **chunk_meta,
                "taxonomy_topic_id": topic.id,
                "taxonomy_topic_name": topic.name,
                "taxonomy_cleaning_policy": "chunk-as-evidence",
            }
        db.add(
            DocumentChunk(
                document_id=document.id,
                knowledge_point_id=point.id,
                chunk_index=index,
                content=safe_content,
                keywords=_extract_keywords(title, section.text),
                page_no=section.page_no,
                slide_no=section.slide_no,
                section_title=_clean_text(section.section_title)[:255],
                token_count=len(safe_content),
                embedding=embed_text(safe_content),
                retrieval_weight=1.0,
                extra_meta=chunk_meta,
            )
        )
    db.flush()
    return len(chunks)


def _ensure_point(
    db: Session,
    chapter: CourseChapter,
    point_cache: dict[str, KnowledgePoint],
    name: str,
    text_value: str,
    source_title: str,
    *,
    prerequisites: Optional[list[str]] = None,
    tags: Optional[list[str]] = None,
    difficulty: Optional[str] = None,
    description: str = "",
) -> KnowledgePoint:
    safe_name = _normalize_point_name(name)[:128] or _normalize_point_name(source_title)[:128] or "imported topic"
    cache_key = _point_cache_key(safe_name)
    point = point_cache.get(cache_key)
    safe_description = _clean_text(description or text_value)[:420]
    source_category = _clean_text(chapter.title)[:80] or "knowledge_base"
    merged_tags = _merge_tags(["imported", "taxonomy_topic", source_category, _clean_text(source_title)[:80]], tags or [])
    merged_prerequisites = [_normalize_point_name(item) for item in (prerequisites or []) if _normalize_point_name(item)]
    if point is None:
        point = KnowledgePoint(
            chapter_id=chapter.id,
            name=safe_name,
            description=safe_description,
            prerequisites=merged_prerequisites,
            tags=merged_tags,
            difficulty=difficulty or _infer_difficulty(safe_name, text_value),
        )
        db.add(point)
        db.flush()
        point_cache[cache_key] = point
    else:
        if safe_description and (len(point.description or "") < 80 or "taxonomy_topic" not in (point.tags or [])):
            point.description = safe_description
        point.tags = _merge_tags(list(point.tags or []), merged_tags)
        point.prerequisites = _merge_tags(list(point.prerequisites or []), merged_prerequisites)
        if difficulty:
            point.difficulty = difficulty
    return point


def _build_import_taxonomy(
    title: str,
    parsed: ParsedKnowledgeFile,
    chunks: list[ParsedSection],
) -> ImportedKnowledgeTaxonomy:
    local = _heuristic_import_taxonomy(title, parsed, chunks)
    try:
        validate_qwen_config()
        taxonomy = qwen_chat_json(
            "You extract a compact learning-topic taxonomy from uploaded course material. Return valid JSON only.",
            _taxonomy_prompt(title, parsed, chunks, local),
            ImportedKnowledgeTaxonomy,
        )
        return _normalize_import_taxonomy(taxonomy, title, chunks, local)
    except Exception:
        return local


def _taxonomy_prompt(
    title: str,
    parsed: ParsedKnowledgeFile,
    chunks: list[ParsedSection],
    local: ImportedKnowledgeTaxonomy,
) -> str:
    section_count = len(parsed.sections)
    chunk_count = len(chunks)
    topic_min, topic_max = _topic_bounds(chunk_count, section_count, parsed.doc_type)
    source_outline = _source_outline(parsed.sections, limit=18)
    sample_chunks = _sample_chunks_for_taxonomy(chunks, limit=14)
    local_seed = [
        {
            "name": topic.name,
            "description": topic.description,
            "source_cues": topic.source_cues[:4],
        }
        for topic in local.topics
    ]
    schema_json = json.dumps(ImportedKnowledgeTaxonomy.model_json_schema(mode="validation"), ensure_ascii=False)
    local_seed_json = json.dumps(local_seed, ensure_ascii=False)
    return "\n".join(
        [
            "请把上传资料抽象成类似 os-taxonomy 的学习主题图谱，而不是把 RAG chunk 当作知识点。",
            f"资料标题：{title}",
            f"资料类型：{parsed.doc_type}",
            f"原始 section 数：{section_count}",
            f"RAG chunk 数：{chunk_count}",
            "",
            "主题数量要求：",
            "- 只保留真正值得学习的小球 topic。",
            f"- 生成 {topic_min}-{topic_max} 个 topic。",
            "- 一节课的教案或 PPT 通常只生成 3 个左右核心知识点；只有内容跨度很大时才增加。",
            "- 不要按 PPT 页码顺序机械生成，不要把每页标题、目录页、参考页、封面页当作 topic。",
            "- 每个 topic 必须不重不漏：名称互不重复，描述覆盖资料重点，但不能把细节碎片都列成 topic。",
            "",
            "topic 字段参考：type, subject, domain, name, description, evidence, assessment_prompt, source_cues。",
            "dependency 字段参考：topic_name, prerequisite_name, strength, reason；必须形成 DAG。",
            "",
            "资料大纲：",
            source_outline,
            "",
            "代表性片段：",
            sample_chunks,
            "",
            "本地初步候选，可合并、改名或删减：",
            local_seed_json,
            "",
            "输出 JSON schema:",
            schema_json,
        ]
    )



def _normalize_import_taxonomy(
    taxonomy: ImportedKnowledgeTaxonomy,
    title: str,
    chunks: list[ParsedSection],
    fallback: ImportedKnowledgeTaxonomy,
) -> ImportedKnowledgeTaxonomy:
    topic_min, topic_max = _topic_bounds(len(chunks), len(chunks), "")
    topics: list[ImportedKnowledgeTopic] = []
    seen: set[str] = set()
    for topic in taxonomy.topics:
        name = _normalize_point_name(topic.name)
        if not name:
            continue
        key = _point_cache_key(name)
        if key in seen:
            continue
        seen.add(key)
        topics.append(
            ImportedKnowledgeTopic(
                id=_clean_text(topic.id)[:64] or _taxonomy_topic_id(name),
                type=topic.type if topic.type in {"CONCEPTUAL", "PROCEDURAL", "REPRESENTATIONAL", "LANGUAGE", "META"} else "CONCEPTUAL",
                subject=_clean_text(topic.subject)[:80] or _clean_text(title)[:80],
                domain=_clean_text(topic.domain)[:80],
                name=name,
                description=_clean_text(topic.description)[:420] or name,
                evidence=[_clean_text(item)[:220] for item in topic.evidence if _clean_text(item)][:4],
                assessment_prompt=_clean_text(topic.assessment_prompt)[:260],
                standards=[_clean_text(item)[:120] for item in topic.standards if _clean_text(item)][:6],
                source_cues=[_clean_text(item)[:80] for item in topic.source_cues if _clean_text(item)][:8],
            )
        )
        if len(topics) >= topic_max:
            break
    if len(topics) < topic_min:
        for topic in fallback.topics:
            key = _point_cache_key(topic.name)
            if key in seen:
                continue
            seen.add(key)
            topics.append(topic)
            if len(topics) >= topic_min:
                break
    if not topics:
        topics = fallback.topics

    names: dict[str, str] = {}
    for topic in topics:
        names[_point_cache_key(topic.name)] = topic.name
        names[_point_cache_key(topic.id)] = topic.name
        names[_point_cache_key(_taxonomy_topic_id(topic.name))] = topic.name
    dependencies: list[ImportedKnowledgeDependency] = []
    seen_edges: set[tuple[str, str]] = set()
    for dependency in taxonomy.dependencies:
        topic_name = names.get(_point_cache_key(dependency.topic_name))
        prereq_name = names.get(_point_cache_key(dependency.prerequisite_name))
        if not topic_name or not prereq_name or topic_name == prereq_name:
            continue
        edge_key = (_point_cache_key(prereq_name), _point_cache_key(topic_name))
        if edge_key in seen_edges:
            continue
        seen_edges.add(edge_key)
        dependencies.append(
            ImportedKnowledgeDependency(
                topic_name=topic_name,
                prerequisite_name=prereq_name,
                strength=dependency.strength if dependency.strength in {"hard", "soft"} else "soft",
                reason=_clean_text(dependency.reason)[:180],
            )
        )
    if not dependencies:
        dependencies = _linear_taxonomy_dependencies(topics)
    return ImportedKnowledgeTaxonomy(topics=topics, dependencies=dependencies[: max(0, len(topics) * 2)])


def _heuristic_import_taxonomy(
    title: str,
    parsed: ParsedKnowledgeFile,
    chunks: list[ParsedSection],
) -> ImportedKnowledgeTaxonomy:
    candidates: list[tuple[str, str, list[str], list[str], int]] = []
    section_groups: dict[str, list[ParsedSection]] = {}
    for section in chunks:
        name = _infer_point_name(title, section)
        if _is_low_value_topic_name(name):
            continue
        key = _point_cache_key(name)
        if not key:
            continue
        section_groups.setdefault(key, []).append(section)

    for sections in section_groups.values():
        first = sections[0]
        name = _infer_point_name(title, first)
        text = _normalize_text("\n".join(section.text for section in sections))
        candidates.append((name, text[:420], _extract_keywords(name, text), _section_cues(sections), len(sections)))

    if not candidates:
        outline = _meaningful_sections(chunks)
        for index, section in enumerate(outline[:8], start=1):
            name = _infer_point_name(title, section) or f"{_normalize_point_name(title)}重点{index}"
            text = _normalize_text(section.text)
            candidates.append((name, text[:420], _extract_keywords(name, text), _section_cues([section]), 1))

    candidates = _merge_topic_candidates(candidates)
    topic_min, topic_max = _topic_bounds(len(chunks), len(parsed.sections), parsed.doc_type)
    ranked = sorted(candidates, key=lambda item: (-item[4], len(item[0]), item[0]))[:topic_max]
    if len(ranked) < topic_min and candidates:
        ranked = candidates[:topic_min]

    topics = [
        ImportedKnowledgeTopic(
            id=_taxonomy_topic_id(name),
            type=_topic_type_from_text(name, description),
            subject=_clean_text(title)[:80],
            domain=_clean_text(parsed.metadata.get("extension") if isinstance(parsed.metadata, dict) else "")[:80],
            name=_normalize_point_name(name),
            description=_clean_text(description)[:360] or _normalize_point_name(name),
            evidence=_topic_evidence(description, keywords),
            assessment_prompt=f"能否说明“{_normalize_point_name(name)}”的核心含义，并结合资料中的例子解释它为什么重要？",
            standards=[],
            source_cues=[*_extract_keywords(name, description)[:4], *cues][:8],
        )
        for name, description, keywords, cues, _count in ranked
        if _normalize_point_name(name)
    ]
    if not topics:
        topics = [
            ImportedKnowledgeTopic(
                id=_taxonomy_topic_id(title),
                type="CONCEPTUAL",
                subject=_clean_text(title)[:80],
                domain="",
                name=_normalize_point_name(title) or "导入资料核心概念",
                description=_summarize_sections(parsed.sections)[:360],
                evidence=["能够概括资料主题", "能够指出关键材料依据"],
                assessment_prompt="能否用自己的话概括这份资料的核心知识？",
                standards=[],
                source_cues=_extract_keywords(title, _summarize_sections(parsed.sections))[:6],
            )
        ]
    return ImportedKnowledgeTaxonomy(topics=topics, dependencies=_linear_taxonomy_dependencies(topics))


def _ensure_taxonomy_points(
    db: Session,
    chapter: CourseChapter,
    point_cache: dict[str, KnowledgePoint],
    taxonomy: ImportedKnowledgeTaxonomy,
    source_title: str,
    chunks: list[ParsedSection],
) -> list[tuple[ImportedKnowledgeTopic, KnowledgePoint]]:
    result: list[tuple[ImportedKnowledgeTopic, KnowledgePoint]] = []
    prereq_by_topic: dict[str, list[str]] = {}
    for dependency in taxonomy.dependencies:
        prereq_by_topic.setdefault(_point_cache_key(dependency.topic_name), []).append(dependency.prerequisite_name)
    corpus = _normalize_text("\n".join(section.text for section in chunks))
    for topic in taxonomy.topics:
        point = _ensure_point(
            db,
            chapter,
            point_cache,
            topic.name,
            topic.description or corpus,
            source_title,
            prerequisites=prereq_by_topic.get(_point_cache_key(topic.name), []),
            tags=_topic_tags(topic, source_title),
            difficulty=_topic_difficulty(topic),
            description=topic.description,
        )
        result.append((topic, point))
    if result:
        return result
    fallback = ImportedKnowledgeTopic(name=_normalize_point_name(source_title) or "导入资料核心概念", description=corpus[:360])
    return [(fallback, _ensure_point(db, chapter, point_cache, fallback.name, corpus, source_title))]


def _match_chunk_to_topic(
    section: ParsedSection,
    topic_points: list[tuple[ImportedKnowledgeTopic, KnowledgePoint]],
) -> KnowledgePoint:
    if not topic_points:
        raise KnowledgeIngestionError("知识点归并失败：没有可用 topic。", 500)
    section_text = _normalize_text(f"{section.section_title}\n{section.text}")
    section_terms = set(_keyword_terms(section_text))
    best_score = -1.0
    best = topic_points[0][1]
    for topic, point in topic_points:
        topic_text = " ".join([topic.name, topic.description, " ".join(topic.source_cues), " ".join(topic.evidence)])
        topic_terms = set(_keyword_terms(topic_text))
        overlap = len(section_terms & topic_terms)
        title_hit = 4 if _point_cache_key(topic.name) and _point_cache_key(topic.name) in _point_cache_key(section.section_title) else 0
        cue_hits = sum(1 for cue in topic.source_cues if cue and cue in section_text)
        score = overlap + title_hit + cue_hits * 2
        if score > best_score:
            best_score = score
            best = point
    return best


def _taxonomy_parse_meta(taxonomy: ImportedKnowledgeTaxonomy) -> dict:
    return {
        "taxonomy_topic_count": len(taxonomy.topics),
        "taxonomy_topics": [
            {
                "name": topic.name,
                "id": topic.id,
                "type": topic.type,
                "subject": topic.subject,
                "domain": topic.domain,
                "evidence": topic.evidence[:4],
                "assessment_prompt": topic.assessment_prompt,
                "assessmentPrompt": topic.assessment_prompt,
                "standards": topic.standards[:6],
            }
            for topic in taxonomy.topics
        ],
        "taxonomy_dependencies": [
            dependency.model_dump(mode="json")
            for dependency in taxonomy.dependencies
        ],
    }


def _taxonomy_topic_id(name: str) -> str:
    digest = hashlib.sha1(_point_cache_key(name).encode("utf-8")).hexdigest()[:12]
    return f"mt_{digest}"


def _topic_bounds(chunk_count: int, section_count: int, doc_type: str) -> tuple[int, int]:
    if chunk_count <= 1:
        return 1, 3
    if doc_type in {"ppt", "pptx"}:
        if section_count <= 12:
            return 2, 4
        if section_count <= 28:
            return 3, 6
        return 4, 8
    if chunk_count <= 6:
        return 2, 5
    if chunk_count <= 20:
        return 3, 7
    return 4, 10


def _source_outline(sections: list[ParsedSection], *, limit: int) -> str:
    lines: list[str] = []
    for index, section in enumerate(sections[:limit], start=1):
        title = _clean_text(section.section_title) or f"section {index}"
        text = _compact_text_for_prompt(section.text, 260)
        location = ""
        if section.slide_no:
            location = f"slide {section.slide_no}"
        elif section.page_no:
            location = f"page {section.page_no}"
        lines.append(f"[{index}] {location} {title}\n{text}")
    return "\n\n".join(lines)


def _sample_chunks_for_taxonomy(chunks: list[ParsedSection], *, limit: int) -> str:
    if not chunks:
        return ""
    if len(chunks) <= limit:
        selected = list(enumerate(chunks, start=1))
    else:
        step = max(1, len(chunks) // limit)
        selected = [(index + 1, chunks[index]) for index in range(0, len(chunks), step)[:limit]]
    return "\n\n".join(
        f"[chunk {index}] {section.section_title}\n{_compact_text_for_prompt(section.text, 420)}"
        for index, section in selected
    )


def _meaningful_sections(sections: list[ParsedSection]) -> list[ParsedSection]:
    return [
        section
        for section in sections
        if len(_normalize_text(section.text)) >= 80
        and not _is_low_value_topic_name(_infer_point_name("", section))
    ]


def _merge_topic_candidates(
    candidates: list[tuple[str, str, list[str], list[str], int]],
) -> list[tuple[str, str, list[str], list[str], int]]:
    merged: list[tuple[str, str, list[str], list[str], int]] = []
    for name, description, keywords, cues, count in candidates:
        normalized = _normalize_point_name(name)
        if not normalized:
            continue
        target_index = -1
        current_terms = set(_keyword_terms(f"{normalized} {description} {' '.join(keywords)}"))
        for index, (existing_name, existing_desc, existing_keywords, _existing_cues, _existing_count) in enumerate(merged):
            existing_terms = set(_keyword_terms(f"{existing_name} {existing_desc} {' '.join(existing_keywords)}"))
            if _point_cache_key(existing_name) == _point_cache_key(normalized) or len(current_terms & existing_terms) >= 3:
                target_index = index
                break
        if target_index < 0:
            merged.append((normalized, description, keywords[:12], cues[:8], count))
            continue
        existing_name, existing_desc, existing_keywords, existing_cues, existing_count = merged[target_index]
        best_name = existing_name if len(existing_name) <= len(normalized) else normalized
        merged[target_index] = (
            best_name,
            existing_desc if len(existing_desc) >= len(description) else description,
            _merge_tags(existing_keywords, keywords)[:12],
            _merge_tags(existing_cues, cues)[:8],
            existing_count + count,
        )
    return merged


def _section_cues(sections: list[ParsedSection]) -> list[str]:
    cues: list[str] = []
    for section in sections[:8]:
        if section.section_title:
            cues.append(_clean_text(section.section_title)[:80])
        if section.page_no:
            cues.append(f"page:{section.page_no}")
        if section.slide_no:
            cues.append(f"slide:{section.slide_no}")
    return _merge_tags([], cues)


def _topic_evidence(description: str, keywords: list[str]) -> list[str]:
    evidence = []
    compact = _compact_text_for_prompt(description, 140)
    if compact:
        evidence.append(compact)
    evidence.extend(f"能解释关键词：{keyword}" for keyword in keywords[:3])
    return evidence[:4] or ["能够说明该知识点的核心含义"]


def _topic_tags(topic: ImportedKnowledgeTopic, source_title: str) -> list[str]:
    return _merge_tags(
        ["taxonomy_topic", topic.type.lower(), topic.subject, topic.domain, _clean_text(source_title)[:80]],
        topic.source_cues,
    )[:12]


def _topic_difficulty(topic: ImportedKnowledgeTopic) -> str:
    raw = f"{topic.name} {topic.description} {' '.join(topic.evidence)}".lower()
    if any(term in raw for term in ["证明", "推导", "实验设计", "优化", "复杂", "advanced", "derive"]):
        return "hard"
    if topic.type in {"PROCEDURAL", "REPRESENTATIONAL"}:
        return "medium"
    if any(term in raw for term in ["定义", "概念", "入门", "基础", "overview", "intro"]):
        return "easy"
    return "medium"


def _topic_type_from_text(name: str, description: str) -> str:
    raw = f"{name} {description}".lower()
    if any(term in raw for term in ["步骤", "流程", "方法", "实验", "操作", "procedure", "workflow"]):
        return "PROCEDURAL"
    if any(term in raw for term in ["图", "模型", "结构", "表示", "公式", "diagram", "representation"]):
        return "REPRESENTATIONAL"
    if any(term in raw for term in ["术语", "定义", "语言", "vocabulary", "term"]):
        return "LANGUAGE"
    if any(term in raw for term in ["反思", "评价", "规范", "meta"]):
        return "META"
    return "CONCEPTUAL"


def _linear_taxonomy_dependencies(topics: list[ImportedKnowledgeTopic]) -> list[ImportedKnowledgeDependency]:
    dependencies: list[ImportedKnowledgeDependency] = []
    for previous, current in zip(topics, topics[1:]):
        dependencies.append(
            ImportedKnowledgeDependency(
                topic_name=current.name,
                prerequisite_name=previous.name,
                strength="soft",
                reason="Inferred from the compact learning order of the uploaded material.",
            )
        )
    return dependencies


def _keyword_terms(value: str) -> list[str]:
    text_value = _clean_text(value).lower()
    ascii_tokens = re.findall(r"[a-z0-9_+\-]{2,}", text_value)
    cjk_tokens = re.findall(r"[\u4e00-\u9fff]{2,8}", text_value)
    terms: list[str] = []
    for token in [*ascii_tokens, *cjk_tokens]:
        token = token.strip("_-+")
        if 2 <= len(token) <= 24 and token not in terms:
            terms.append(token)
    return terms[:80]


def _merge_tags(left: list[str], right: list[str]) -> list[str]:
    result: list[str] = []
    for item in [*left, *right]:
        value = _clean_text(item)
        if value and value not in result:
            result.append(value)
    return result[:16]


def _is_low_value_topic_name(value: str) -> bool:
    normalized = _normalize_point_name(value).lower()
    if not normalized:
        return True
    low_value = {
        "目录",
        "contents",
        "参考文献",
        "references",
        "谢谢",
        "thank you",
        "课程介绍",
        "导入资料",
    }
    if normalized in low_value:
        return True
    return bool(re.fullmatch(r"(slide|page|section|chapter)?\s*\d+", normalized))


def _compact_text_for_prompt(value: object, limit: int) -> str:
    return re.sub(r"\s+", " ", _clean_text(value)).strip()[:limit]


def _chunk_sections(sections: list[ParsedSection]) -> list[ParsedSection]:
    settings = get_settings()
    size = max(400, settings.knowledge_chunk_chars)
    overlap = max(0, min(settings.knowledge_chunk_overlap, size // 2))
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
    except ImportError as exc:
        raise KnowledgeIngestionError("缺少 RAG 文档切分依赖 langchain-text-splitters，请安装 backend/requirements.txt。", 500) from exc

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", "。", "！", "？", ";", "；", " ", ""],
    )
    chunks: list[ParsedSection] = []
    for section in sections:
        text_value = _normalize_text(section.text)
        for part_index, part in enumerate(splitter.split_text(text_value), start=1):
            part = _normalize_text(part)
            if not part:
                continue
            section_title = section.section_title if len(text_value) <= size else f"{section.section_title} #{part_index}".strip()
            chunks.append(
                ParsedSection(
                    text=part,
                    section_title=section_title,
                    page_no=section.page_no,
                    slide_no=section.slide_no,
                    metadata=section.metadata,
                )
            )
    return chunks


def search_knowledge_enhanced(db: Session, query: str, limit: int = 8) -> list[dict]:
    vector = embed_text(query)
    terms = [term.strip() for term in query.replace("，", " ").replace("、", " ").replace("/", " ").split() if term.strip()]
    keyword_filters = " OR ".join(
        [
            "dc.content ILIKE :pattern",
            "kp.name ILIKE :pattern",
            "kd.title ILIKE :pattern",
            "kd.summary ILIKE :pattern",
        ]
    )
    pattern = f"%{terms[0] if terms else query}%"
    sql = text(
        f"""
        SELECT
            dc.id AS chunk_id,
            kd.title AS document_title,
            kd.doc_type AS document_type,
            COALESCE(kp.name, '课程资料') AS knowledge_point,
            dc.content AS content,
            kd.source_uri AS source_uri,
            dc.keywords AS keywords,
            dc.page_no AS page_no,
            dc.slide_no AS slide_no,
            dc.section_title AS section_title,
            CASE WHEN dc.embedding IS NULL THEN NULL ELSE dc.embedding <=> CAST(:embedding AS vector) END AS distance,
            CASE WHEN ({keyword_filters}) THEN 1 ELSE 0 END AS keyword_hit
        FROM document_chunks dc
        JOIN knowledge_documents kd ON dc.document_id = kd.id
        LEFT JOIN knowledge_points kp ON dc.knowledge_point_id = kp.id
        ORDER BY
            keyword_hit DESC,
            CASE WHEN dc.embedding IS NULL THEN 1 ELSE 0 END ASC,
            distance ASC NULLS LAST,
            dc.retrieval_weight DESC,
            dc.id DESC
        LIMIT :limit
        """
    )
    rows = db.execute(sql, {"embedding": _vector_literal(vector), "pattern": pattern, "limit": limit}).mappings().all()
    return [dict(row) for row in rows]


def build_rag_context(db: Session, query: str, limit: int = 8) -> str:
    hits = search_knowledge_enhanced(db, query, limit=limit)
    if not hits:
        return "未检索到课程知识库资料。请要求用户补充课件、论文或实验资料。"
    lines = []
    for index, hit in enumerate(hits, start=1):
        loc = []
        if hit.get("page_no"):
            loc.append(f"page {hit['page_no']}")
        if hit.get("slide_no"):
            loc.append(f"slide {hit['slide_no']}")
        location = f" ({', '.join(loc)})" if loc else ""
        lines.append(
            f"[{index}] {hit['knowledge_point']} / {hit['document_title']}{location} / {hit['document_type']}\n"
            f"source={hit['source_uri']}\n"
            f"{hit['content'][:1200]}"
        )
    return "\n\n".join(lines)


def embed_text(text_value: str) -> list[float]:
    settings = get_settings()
    if settings.knowledge_embedding_provider.lower() != "qwen":
        raise KnowledgeIngestionError(f"不支持的 embedding provider：{settings.knowledge_embedding_provider}", 501)
    validate_qwen_config()
    import urllib.request
    import urllib.error

    payload = {
        "model": settings.knowledge_embedding_model,
        "input": text_value[:8192],
        "encoding_format": "float",
    }
    data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
    request = urllib.request.Request(
        f"{settings.qwen_base_url.rstrip('/')}/embeddings",
        data=data,
        headers={
            "Authorization": f"Bearer {settings.qwen_api_key}",
            "Content-Type": "application/json",
        },
        method="POST",
    )
    try:
        with urllib.request.urlopen(request, timeout=settings.qwen_timeout_seconds) as response:
            body = json.loads(response.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="ignore")
        raise LLMResponseError(f"Qwen Embedding 接口返回错误：{exc.code} {detail}") from exc
    except urllib.error.URLError as exc:
        raise LLMResponseError(f"无法连接 Qwen Embedding 接口：{exc.reason}") from exc
    try:
        vector = body["data"][0]["embedding"]
    except (KeyError, IndexError, TypeError) as exc:
        raise LLMResponseError("Qwen Embedding 响应缺少 data[0].embedding") from exc
    if len(vector) != settings.knowledge_embedding_dim:
        raise LLMResponseError(f"Embedding 维度不匹配：返回 {len(vector)}，配置 {settings.knowledge_embedding_dim}")
    return vector


def _infer_point_name(title: str, section: ParsedSection) -> str:
    section_title = _normalize_point_name(section.section_title)
    if section_title and not _is_generic_section_title(section_title):
        return section_title
    for line in section.text.splitlines()[:8]:
        candidate = _normalize_point_name(line)
        if candidate and not _is_low_value_topic_name(candidate):
            return candidate
    return _normalize_point_name(title) or "imported topic"


def _normalize_point_name(value: object) -> str:
    text_value = _clean_text(value)
    text_value = re.sub(r"^第?\s*\d+\s*[章节讲页幻灯片、.\-_\s]*", "", text_value, flags=re.IGNORECASE)
    text_value = re.sub(r"^(chapter|section|slide|page)\s*\d*[:：.\-\s]*", "", text_value, flags=re.IGNORECASE)
    text_value = re.sub(r"^[#*\-\s·•]+", "", text_value)
    text_value = re.sub(r"\s+", " ", text_value).strip(" ：:;；，,。.")
    if not text_value:
        return ""
    if len(text_value) > 48:
        sentence = re.split(r"[。！？!?；;：:\n]", text_value, maxsplit=1)[0].strip()
        if 2 <= len(sentence) <= 48:
            text_value = sentence
    if len(text_value) > 48:
        keywords = _extract_keywords("", text_value)
        if keywords:
            text_value = " / ".join(keywords[:3])
    return text_value[:48]


def _point_cache_key(value: object) -> str:
    text_value = _normalize_point_name(value).lower()
    return re.sub(r"[\s\-_:：/\\（）()【】\[\]·•]+", "", text_value)


def _is_generic_section_title(value: str) -> bool:
    text_value = _clean_text(value).lower()
    if not text_value:
        return True
    return bool(
        re.fullmatch(r"(第\s*)?\d+\s*(页|頁|幻灯片|slide|page|ocr)?", text_value, flags=re.IGNORECASE)
        or re.fullmatch(r"(slide|page)\s*\d+", text_value, flags=re.IGNORECASE)
    )


def _infer_difficulty(name: str, text_value: str) -> str:
    raw = f"{name} {text_value[:420]}".lower()
    if any(term in raw for term in ["证明", "推导", "复杂", "高级", "优化", "hard", "advanced"]):
        return "hard"
    if any(term in raw for term in ["定义", "概念", "基础", "入门", "介绍", "basic", "intro"]):
        return "easy"
    return "medium"


def _extract_keywords(title: str, text_value: str) -> list[str]:
    raw = f"{title} {text_value[:500]}"
    terms = []
    for term in raw.replace("，", " ").replace("。", " ").replace("、", " ").replace("/", " ").split():
        cleaned = term.strip(" ：:;；,.，。()[]【】")
        if 2 <= len(cleaned) <= 24 and cleaned not in terms:
            terms.append(cleaned)
        if len(terms) >= 12:
            break
    return terms



def _summarize_sections(sections: list[ParsedSection]) -> str:
    text_value = "\n".join(section.text for section in sections)
    return _normalize_text(text_value)[:1000]


def _normalize_text(text_value: str) -> str:
    cleaned = _clean_text(text_value)
    return "\n".join(line.strip() for line in cleaned.splitlines() if line.strip())


def _clean_text(value: object) -> str:
    text_value = "" if value is None else str(value)
    text_value = text_value.encode("utf-8", errors="ignore").decode("utf-8", errors="ignore")
    text_value = CONTROL_CHARS_RE.sub("", text_value)
    return text_value.strip()


def _clean_metadata(value: object) -> object:
    if isinstance(value, dict):
        return {_clean_text(key): _clean_metadata(item) for key, item in value.items() if _clean_text(key)}
    if isinstance(value, list):
        return [_clean_metadata(item) for item in value]
    if isinstance(value, tuple):
        return [_clean_metadata(item) for item in value]
    if isinstance(value, str):
        return _clean_text(value)
    return value


def _safe_filename(filename: str) -> str:
    recovered = _recover_zip_mojibake(filename)
    value = _clean_text(Path(recovered).name)
    return value or "upload.bin"


def _recover_zip_mojibake(value: str) -> str:
    text_value = _clean_text(value)
    if not text_value:
        return text_value
    try:
        recovered = text_value.encode("cp437").decode("gb18030")
    except (UnicodeEncodeError, UnicodeDecodeError):
        return text_value
    return recovered if any("\u4e00" <= char <= "\u9fff" for char in recovered) else text_value


def _job_storage_dir(job_id: int) -> Path:
    return Path(get_settings().knowledge_upload_dir) / f"job-{job_id}"


def _job_storage_bytes(job_id: int) -> int:
    job_dir = _job_storage_dir(job_id)
    if not job_dir.exists():
        return 0
    return sum(path.stat().st_size for path in job_dir.rglob("*") if path.is_file())


def _documents_for_job(db: Session, job_id: int) -> list[KnowledgeDocument]:
    documents = db.scalars(select(KnowledgeDocument).where(KnowledgeDocument.file_path.ilike("%job-%"))).all()
    return [document for document in documents if _job_id_from_document_path(document.file_path) == job_id]


def _delete_job_storage_dir(job_id: int) -> None:
    job_dir = _job_storage_dir(job_id)
    if job_dir.exists():
        shutil.rmtree(job_dir, ignore_errors=True)


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _vector_literal(vector: list[float]) -> str:
    return "[" + ",".join(str(float(value)) for value in vector) + "]"
